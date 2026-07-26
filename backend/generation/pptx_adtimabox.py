"""
AdtimaBox-branded PPTX generator.
Slide types + extraction schema: agents/wireframe_designer_agent/SKILL.md (same source
html_deck.py loads) — all 14 types (cover, closing, agenda, split, highlight, value, flow,
touchpoints, compliance, tier, roi, timeline, checklist, screen) must stay handled here in
sync with html_deck.py's `_render_slide`, or a slide silently falls back to `_render_value`
and renders title-only with a blank body.
"""

from __future__ import annotations

import asyncio
import json
import os
from functools import partial
from typing import Any

# Brand palette — RGB (0-255)
_C = {
    "orange":  (246, 80, 9),
    "orange2": (232, 74, 26),
    "teal":    (15, 155, 142),
    "ink":     (29, 29, 31),
    "white":   (255, 255, 255),
    "cream":   (251, 248, 245),
    "line":    (236, 230, 225),
    "gray":    (107, 107, 112),
    "gray_lt": (154, 154, 160),
    "purple":  (91, 79, 196),
    "gold":    (200, 147, 43),
}

SLIDE_W = 10.0     # inches
SLIDE_H = 5.625    # inches
FONT = "Calibri"

# Layout constants (align with adtimabox-deck-pptx.skill)
PAD_X = 0.67   # left + right padding
PAD_T = 0.44   # top padding (TOPBAR_H + gap)
TOPBAR_H = 0.18
STAT_BAR_H = 0.55
CONTENT_W = SLIDE_W - PAD_X * 2
BODY_TOP = PAD_T + TOPBAR_H + 0.18   # below topbar
BODY_H = SLIDE_H - BODY_TOP - STAT_BAR_H - 0.10

# Value slide right-stat-column (mirrors HTML rs-card-primary/secondary)
RC_W = 2.10                          # right col width
RC_X = SLIDE_W - PAD_X - RC_W       # = 7.23" — right col left edge
LEFT_W_V = RC_X - PAD_X - 0.38      # = 6.18" — left body width (gap 0.38")


class AdtimaBoxPPTXGenerator:
    """Generates AdtimaBox-branded PPTX from proposal markdown."""

    async def generate(self, proposal_text: str, brief: dict, output_path: str, skill_spec: str = "") -> dict:
        try:
            from pptx import Presentation  # noqa
        except ImportError:
            return {"status": "error", "error": "python-pptx not installed"}

        try:
            # Share extraction with html_deck (same schema, retry logic included)
            from generation.html_deck import HTMLDeckGenerator
            helper = HTMLDeckGenerator()
            slides_data = await helper._extract_slides_with_retry(proposal_text, brief)
            if not slides_data:
                slides_data = self._fallback_slides(brief)
        except Exception as e:
            print(f"[PPTX] Extraction failed ({e}), using fallback")
            slides_data = self._fallback_slides(brief)

        try:
            prs = self._build_pptx(slides_data)
            os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
            prs.save(output_path)
            return {"status": "success", "file_path": output_path, "slide_count": len(slides_data)}
        except Exception as e:
            return {"status": "error", "error": str(e)}

    def _fallback_slides(self, brief: dict) -> list[dict]:
        # Delegate to HTMLDeckGenerator so both generators stay in sync
        from generation.html_deck import HTMLDeckGenerator
        helper = HTMLDeckGenerator()
        return helper._ensure_required_slides([], brief)

    def _build_pptx(self, slides: list[dict]):
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        prs.slide_width = Inches(SLIDE_W)
        prs.slide_height = Inches(SLIDE_H)

        dispatch = {
            "cover":       self._render_cover,
            "closing":     self._render_closing,
            "agenda":      self._render_agenda,
            "split":       self._render_split,
            "highlight":   self._render_highlight,
            "value":       self._render_value,
            "flow":        self._render_flow,
            "touchpoints": self._render_touchpoints,
            "compliance":  self._render_compliance,
            "tier":        self._render_tier,
            "roi":         self._render_roi,
            "timeline":    self._render_timeline,
            "checklist":   self._render_checklist,
            "screen":      self._render_screen,
        }
        blank = prs.slide_layouts[6]
        for sd in slides:
            slide = prs.slides.add_slide(blank)
            renderer = dispatch.get(sd.get("type", "value"), self._render_value)
            renderer(slide, sd)
        return prs

    # ── Helpers ──────────────────────────────────────────────────────────

    @staticmethod
    def _rgb(key: str):
        from pptx.dml.color import RGBColor
        r, g, b = _C[key]
        return RGBColor(r, g, b)

    @staticmethod
    def _hex_rgb(hex_str: str):
        from pptx.dml.color import RGBColor
        h = hex_str.lstrip("#").upper()
        if len(h) != 6:
            return RGBColor(29, 29, 31)
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    @staticmethod
    def _round_corners(shape, adj: int = 10000):
        """Make a rectangle shape have rounded corners.
        adj — corner radius as a fraction of the shorter edge × 100000.
        Common values: 10000 = 10% (cards), 20000 = 20% (icon boxes), 50000 = pill.
        Note: spPr lives in the p: (presentation) namespace; prstGeom/avLst/gd in a: (drawingml).
        """
        from lxml import etree
        p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
        a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        sp = shape._element
        spPr = sp.find(f"{{{p_ns}}}spPr")
        if spPr is None:
            return
        prstGeom = spPr.find(f"{{{a_ns}}}prstGeom")
        if prstGeom is None:
            return
        prstGeom.set("prst", "roundRect")
        for old in prstGeom.findall(f"{{{a_ns}}}avLst"):
            prstGeom.remove(old)
        avLst = etree.SubElement(prstGeom, f"{{{a_ns}}}avLst")
        gd = etree.SubElement(avLst, f"{{{a_ns}}}gd")
        gd.set("name", "adj")
        gd.set("fmla", f"val {adj}")

    def _bg(self, slide):
        from pptx.util import Inches
        bg = slide.shapes.add_shape(1, 0, 0, Inches(SLIDE_W), Inches(SLIDE_H))
        bg.fill.solid()
        bg.fill.fore_color.rgb = self._rgb("cream")
        bg.line.fill.background()

    def _topbar(self, slide, eyebrow: str, tier: str = ""):
        from pptx.util import Inches, Pt
        # Small orange accent bar (4px × 16px equivalent in inches ≈ 0.04" × 0.13")
        bar = slide.shapes.add_shape(1, Inches(PAD_X), Inches(PAD_T + 0.03),
                                     Inches(0.04), Inches(0.14))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self._rgb("orange")
        bar.line.fill.background()

        # Eyebrow text
        label = eyebrow.upper()
        if tier:
            label += f"  {tier.upper()}"
        tb = slide.shapes.add_textbox(Inches(PAD_X + 0.10), Inches(PAD_T),
                                      Inches(CONTENT_W * 0.6), Inches(TOPBAR_H))
        tf = tb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = label
        run.font.name = FONT
        run.font.size = Pt(9)
        run.font.bold = True
        run.font.color.rgb = self._rgb("ink")

        # Logo — right aligned: "adtimabox  by Adtima"
        logo_tb = slide.shapes.add_textbox(Inches(SLIDE_W - PAD_X - 1.8), Inches(PAD_T),
                                           Inches(1.8), Inches(TOPBAR_H))
        lf = logo_tb.text_frame
        lf.word_wrap = False
        lp = lf.paragraphs[0]
        from pptx.enum.text import PP_ALIGN
        lp.alignment = PP_ALIGN.RIGHT
        lr = lp.add_run()
        lr.text = "adtimabox"
        lr.font.name = FONT
        lr.font.size = Pt(9)
        lr.font.bold = True
        lr.font.color.rgb = self._rgb("ink")
        lr2 = lp.add_run()
        lr2.text = "  by Adtima"
        lr2.font.name = FONT
        lr2.font.size = Pt(8)
        lr2.font.bold = False
        lr2.font.color.rgb = self._rgb("gray_lt")

    def _stat_bar(self, slide, stats: list, no_line: bool = False):
        from pptx.util import Inches, Pt
        if not stats:
            return
        # Divider line (omit when caller has own footer divider)
        if not no_line:
            div = slide.shapes.add_shape(1, Inches(PAD_X),
                                         Inches(SLIDE_H - STAT_BAR_H - 0.02),
                                         Inches(CONTENT_W), Inches(0.01))
            div.fill.solid()
            div.fill.fore_color.rgb = self._rgb("line")
            div.line.fill.background()

        stat_x = PAD_X
        # 22pt bold text needs ~0.37" for a single line (font_pt * 1.22 / 72) — the old 0.30"
        # box clipped/overlapped even short values. value(0.36) + label(0.16) = 0.52" within
        # STAT_BAR_H 0.55".
        stat_y = SLIDE_H - STAT_BAR_H + 0.02
        col_w = CONTENT_W / min(len(stats), 4)
        for s in stats[:4]:
            self._text(slide, stat_x, stat_y, col_w - 0.1, 0.36,
                       s.get("v", ""), 22, "ink", bold=True, auto_fit=True)
            self._text(slide, stat_x, stat_y + 0.36, col_w - 0.1, 0.16,
                       s.get("l", ""), 8, "gray_lt")
            stat_x += col_w

    @staticmethod
    def _safe_text(text: str) -> str:
        """Strip characters that break OOXML: lone surrogates, variation selectors, and C0/C1 controls."""
        if not text:
            return ""
        result = []
        for ch in str(text):
            cp = ord(ch)
            # Lone surrogates (U+D800–U+DFFF) are illegal in XML
            if 0xD800 <= cp <= 0xDFFF:
                continue
            # Variation selectors (U+FE00–U+FE0F, U+E0100–U+E01EF) that arrive
            # without their base emoji cause lxml to fail
            if 0xFE00 <= cp <= 0xFE0F or 0xE0100 <= cp <= 0xE01EF:
                continue
            # C0 controls except tab/LF/CR
            if cp < 0x20 and ch not in "\t\n\r":
                continue
            result.append(ch)
        return "".join(result)

    @staticmethod
    def _est_lines(text: str, width_in: float, font_pt: float) -> int:
        """Estimate wrapped line count for `text` at `font_pt` inside a `width_in`-wide box.
        Unlike HTML/CSS, python-pptx text boxes don't reflow their container — callers must
        size boxes to fit real (often multi-line Vietnamese) content, or PowerPoint silently
        overflows text on top of whatever shape sits below."""
        import math
        if not text:
            return 0
        avg_char_w_in = font_pt * 0.52 / 72.0
        chars_per_line = max(int(width_in / avg_char_w_in), 1)
        lines = 0
        for para in str(text).split("\n"):
            lines += max(1, math.ceil(len(para) / chars_per_line))
        return lines

    def _est_h(self, text: str, width_in: float, font_pt: float, line_factor: float = 1.22) -> float:
        """Estimated rendered height (inches) of `text` wrapped at `width_in` and `font_pt`."""
        return self._est_lines(text, width_in, font_pt) * font_pt * line_factor / 72.0

    def _fit_list_height(self, items: list, width_in: float, font_pt: float, max_h_in: float,
                          item_gap: float = 0.06, min_item_h: float = 0.22,
                          line_factor: float = 1.22):
        """Greedily keep as many `items` as fit within `max_h_in`, sizing each to its real
        wrapped height instead of a fixed per-row increment. Returns (kept_items, heights,
        dropped_count) — always keeps at least one item so we never render nothing."""
        heights: list[float] = []
        kept: list = []
        total = 0.0
        for it in items:
            text = it.get("text", it) if isinstance(it, dict) else it
            h = max(self._est_h(text, width_in, font_pt, line_factor) + item_gap, min_item_h)
            if total + h > max_h_in and kept:
                break
            kept.append(it)
            heights.append(h)
            total += h
        return kept, heights, len(items) - len(kept)

    def _text(self, slide, left, top, width, height, text, size, color_key,
              bold=False, italic=False, align="LEFT", emoji_font=False, auto_fit=False):
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        aligns = {"LEFT": PP_ALIGN.LEFT, "CENTER": PP_ALIGN.CENTER, "RIGHT": PP_ALIGN.RIGHT}
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = tb.text_frame
        tf.word_wrap = True
        if auto_fit:
            from pptx.enum.text import MSO_AUTO_SIZE
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        p = tf.paragraphs[0]
        p.alignment = aligns.get(align, PP_ALIGN.LEFT)
        run = p.add_run()
        run.text = self._safe_text(text or "")
        # Emoji icons need a font that has emoji glyphs (Segoe UI Emoji on Windows)
        run.font.name = "Segoe UI Emoji" if emoji_font else FONT
        run.font.size = Pt(size)
        run.font.color.rgb = self._rgb(color_key)
        run.font.bold = bold
        run.font.italic = italic
        return tb

    def _text_mixed(self, slide, left, top, width, height, plain: str, bold_text: str, size: int):
        """Two-run paragraph: plain + bold (orange) for headline:{plain,bold}."""
        from pptx.util import Inches, Pt
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = self._safe_text(plain or "")
        r1.font.name = FONT
        r1.font.size = Pt(size)
        r1.font.color.rgb = self._rgb("ink")
        if bold_text:
            r2 = p.add_run()
            r2.text = self._safe_text(bold_text)
            r2.font.name = FONT
            r2.font.size = Pt(size)
            r2.font.color.rgb = self._rgb("orange")
            r2.font.bold = True
        return tb

    # ── Slide renderers ──────────────────────────────────────────────────

    def _dark_bg(self, slide):
        from pptx.util import Inches
        bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(SLIDE_W), Inches(SLIDE_H))
        bg.fill.solid()
        bg.fill.fore_color.rgb = self._rgb("ink")
        bg.line.fill.background()
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(SLIDE_W), Inches(0.05))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self._rgb("orange")
        bar.line.fill.background()

    def _render_cover(self, slide, sd: dict):
        """Dark hero cover slide (always first)."""
        self._dark_bg(slide)
        brand = sd.get("brand", "")
        industry = sd.get("industry", "")
        date = sd.get("date", "")
        track = sd.get("track", "")

        self._text(slide, PAD_X, 1.10, CONTENT_W, 0.30, brand.upper(), 12, "orange", bold=True)
        self._text(slide, PAD_X, 1.48, CONTENT_W, 0.90, "ADTIMABOX PROPOSAL", 32, "white", bold=True)
        self._text(slide, PAD_X, 2.45, CONTENT_W, 0.30,
                   f"Zalo Ecosystem Solutions — {industry or 'Client Pitch'}", 12, "gray_lt")

        meta = [(l, v) for l, v in (("Date", date), ("Industry", industry), ("Track", track)) if v]
        if meta:
            mx = PAD_X
            col_w = 2.3
            meta_y = SLIDE_H - 0.85
            for label, val in meta:
                self._text(slide, mx, meta_y, col_w - 0.2, 0.18, label.upper(), 7.5, "gray_lt")
                self._text(slide, mx, meta_y + 0.20, col_w - 0.2, 0.24, val, 10, "white", bold=True)
                mx += col_w
        self._text(slide, SLIDE_W - PAD_X - 1.2, SLIDE_H - 0.35, 1.2, 0.2, "ADTIMA", 7.5, "gray_lt", align="RIGHT")

    def _render_closing(self, slide, sd: dict):
        """Dark thank-you closing slide (always last)."""
        from pptx.util import Pt
        self._dark_bg(slide)
        brand = sd.get("brand", "")
        date = sd.get("date", "")
        note = sd.get("note", "Confidential. Prices excl. VAT 8%. Valid 30 days from proposal date.")

        from pptx.util import Inches
        tag = slide.shapes.add_shape(1, Inches(PAD_X), Inches(1.10), Inches(2.1), Inches(0.30))
        tag.fill.solid()
        tag.fill.fore_color.rgb = self._rgb("orange")
        tag.line.fill.background()
        self._round_corners(tag, 15000)
        self._text(slide, PAD_X, 1.10, 2.1, 0.30, "ADTIMABOX BY ADTIMA", 8, "white", bold=True, align="CENTER")

        self._text(slide, PAD_X, 1.60, CONTENT_W, 0.65, f"Cảm ơn {brand}", 28, "white", bold=True)
        self._text(slide, PAD_X, 2.30, CONTENT_W, 0.30,
                   "Hẹn gặp lại — chúng tôi sẵn sàng hỗ trợ bạn.", 12, "gray_lt")

        meta_y = SLIDE_H - 1.05
        self._text(slide, PAD_X, meta_y, 2.4, 0.18, "PREPARED BY", 7.5, "gray_lt")
        self._text(slide, PAD_X, meta_y + 0.20, 2.4, 0.24, "AdtimaBox Sales Team", 10, "white", bold=True)
        if date:
            self._text(slide, PAD_X + 2.6, meta_y, 2.0, 0.18, "DATE", 7.5, "gray_lt")
            self._text(slide, PAD_X + 2.6, meta_y + 0.20, 2.0, 0.24, date, 10, "white", bold=True)

        self._text(slide, PAD_X, SLIDE_H - 0.55, CONTENT_W, 0.20, note, 8, "gray_lt")
        self._text(slide, SLIDE_W - PAD_X - 1.2, SLIDE_H - 0.35, 1.2, 0.2, "ADTIMA", 7.5, "gray_lt", align="RIGHT")

    def _render_agenda(self, slide, sd: dict):
        """Numbered grid of section names present in the proposal."""
        from pptx.util import Inches, Pt
        self._bg(slide)
        self._topbar(slide, sd.get("eyebrow", "Agenda"))

        items = (sd.get("items") or [])[:12]
        n = max(len(items), 1)
        cols = 2 if n <= 6 else 3
        rows = (n + cols - 1) // cols
        gap_x, gap_y = 0.16, 0.14
        item_w = (CONTENT_W - gap_x * (cols - 1)) / cols
        item_h = min((BODY_H - gap_y * (rows - 1)) / rows, 0.55)
        block_h = item_h * rows + gap_y * (rows - 1)
        start_y = BODY_TOP + max((BODY_H - block_h) / 2, 0)

        for i, item in enumerate(items):
            r, c = divmod(i, cols)
            x = PAD_X + c * (item_w + gap_x)
            y = start_y + r * (item_h + gap_y)
            card = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(item_w), Inches(item_h))
            card.fill.solid()
            card.fill.fore_color.rgb = self._rgb("cream")
            card.line.color.rgb = self._rgb("line")
            card.line.width = Pt(0.5)
            self._round_corners(card, 10000)

            num_size = min(item_h - 0.16, 0.30)
            num_y = y + (item_h - num_size) / 2
            badge = slide.shapes.add_shape(9, Inches(x + 0.10), Inches(num_y),
                                           Inches(num_size), Inches(num_size))
            badge.fill.solid()
            badge.fill.fore_color.rgb = self._rgb("orange")
            badge.line.fill.background()
            self._text(slide, x + 0.10, num_y, num_size, num_size,
                       str(i + 1), 10, "white", bold=True, align="CENTER")
            self._text(slide, x + 0.10 + num_size + 0.10, y, item_w - num_size - 0.30, item_h,
                       item, 10.5, "ink", bold=True)

    def _render_split(self, slide, sd: dict):
        """Two-panel business-problem slide (AS-IS / TO-BE)."""
        from pptx.util import Inches, Pt
        self._bg(slide)
        self._topbar(slide, sd.get("eyebrow", "Business Problem"))

        left_label = sd.get("left_label", "Current State")
        left_text = sd.get("left_text", "")
        left_pain = sd.get("left_pain", "")
        right_label = sd.get("right_label", "Target State")
        right_text = sd.get("right_text", "")
        gap_text = sd.get("gap", "")
        right_items = (sd.get("right_items") or [])[:4]
        stats = sd.get("stats") or []

        panel_gap = 0.24
        panel_w = (CONTENT_W - panel_gap) / 2
        panel_h = BODY_H
        panel_y = BODY_TOP
        pain_tint = self._hex_rgb("FDEAE0")

        # Left panel — orange accent border
        lp = slide.shapes.add_shape(1, Inches(PAD_X), Inches(panel_y), Inches(panel_w), Inches(panel_h))
        lp.fill.solid()
        lp.fill.fore_color.rgb = self._rgb("cream")
        lp.line.color.rgb = self._rgb("orange")
        lp.line.width = Pt(1.25)
        self._round_corners(lp, 6000)
        self._text(slide, PAD_X + 0.16, panel_y + 0.14, panel_w - 0.32, 0.20,
                   left_label.upper(), 8.5, "orange", bold=True)
        pain_h = 0.42 if left_pain else 0
        self._text(slide, PAD_X + 0.16, panel_y + 0.42, panel_w - 0.32,
                   panel_h - 0.58 - pain_h, left_text, 9.5, "gray")
        if left_pain:
            py = panel_y + panel_h - 0.50
            box = slide.shapes.add_shape(1, Inches(PAD_X + 0.16), Inches(py),
                                         Inches(panel_w - 0.32), Inches(0.38))
            box.fill.solid()
            box.fill.fore_color.rgb = pain_tint
            box.line.fill.background()
            self._round_corners(box, 12000)
            self._text(slide, PAD_X + 0.26, py + 0.04, panel_w - 0.52, 0.30,
                       f'"{left_pain}"', 8.5, "ink", italic=True)

        # Right panel
        rx = PAD_X + panel_w + panel_gap
        rp = slide.shapes.add_shape(1, Inches(rx), Inches(panel_y), Inches(panel_w), Inches(panel_h))
        rp.fill.solid()
        rp.fill.fore_color.rgb = self._rgb("cream")
        rp.line.color.rgb = self._rgb("line")
        rp.line.width = Pt(0.5)
        self._round_corners(rp, 6000)
        self._text(slide, rx + 0.16, panel_y + 0.14, panel_w - 0.32, 0.20,
                   right_label.upper(), 8.5, "orange", bold=True)
        # right_text routinely wraps past the old fixed 0.60"/2-line budget — size it to its
        # real height (min 0.60 to match prior look for short text) so items never overlap it.
        rt_w = panel_w - 0.32
        rt_h = max(self._est_h(right_text, rt_w, 9.5) + 0.04, 0.60)
        self._text(slide, rx + 0.16, panel_y + 0.42, rt_w, rt_h, right_text, 9.5, "gray")
        cy = panel_y + 0.42 + rt_h + 0.04
        gap_h = 0.36 if gap_text else 0
        if gap_text:
            box = slide.shapes.add_shape(1, Inches(rx + 0.16), Inches(cy),
                                         Inches(panel_w - 0.32), Inches(gap_h))
            box.fill.solid()
            box.fill.fore_color.rgb = pain_tint
            box.line.fill.background()
            self._round_corners(box, 12000)
            self._text(slide, rx + 0.26, cy + 0.03, panel_w - 0.52, gap_h - 0.06,
                       gap_text, 8.5, "ink", italic=True)
            cy += gap_h + 0.08
        items_avail = panel_y + panel_h - 0.10 - cy
        kept, heights, more = self._fit_list_height(
            right_items, panel_w - 0.32 - 0.20, 9, max(items_avail, 0), min_item_h=0.24)
        for it, h in zip(kept, heights):
            self._text(slide, rx + 0.16, cy, panel_w - 0.32, h - 0.02, f"→  {it}", 9, "ink")
            cy += h
        if more:
            self._text(slide, rx + 0.16, cy, panel_w - 0.32, 0.20, f"+{more} khác…", 8, "gray_lt", italic=True)

        self._stat_bar(slide, stats)

    def _render_touchpoints(self, slide, sd: dict):
        """Messaging-map table slide."""
        from pptx.util import Inches, Pt
        self._bg(slide)
        self._topbar(slide, sd.get("eyebrow", "Messaging Touchpoints"))
        hl = sd.get("headline", {})
        self._text_mixed(slide, PAD_X, BODY_TOP, CONTENT_W, 0.42,
                         hl.get("plain", ""), hl.get("bold", ""), 18)

        rows = (sd.get("rows") or [])[:8]
        stats = sd.get("stats") or []
        table_top = BODY_TOP + 0.46
        table_h = BODY_H - 0.46
        n_rows = len(rows) + 1  # + header

        if rows:
            headers = ["Trigger", "Message Type", "Channel", "Timing"]
            gfx = slide.shapes.add_table(n_rows, len(headers),
                                         Inches(PAD_X), Inches(table_top),
                                         Inches(CONTENT_W), Inches(table_h))
            table = gfx.table
            col_w = [0.34, 0.24, 0.20, 0.22]
            for i, w in enumerate(col_w):
                table.columns[i].width = Inches(CONTENT_W * w)

            for c, h in enumerate(headers):
                cell = table.cell(0, c)
                cell.text = h
                cell.fill.solid()
                cell.fill.fore_color.rgb = self._rgb("cream")
                run = cell.text_frame.paragraphs[0].runs[0]
                run.font.name = FONT
                run.font.size = Pt(8.5)
                run.font.bold = True
                run.font.color.rgb = self._rgb("gray")

            for r, row in enumerate(rows, start=1):
                values = [row.get("trigger", ""), row.get("message_type", ""),
                         row.get("channel", ""), row.get("timing", "")]
                for c, val in enumerate(values):
                    cell = table.cell(r, c)
                    cell.text = self._safe_text(str(val))
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self._rgb("white")
                    run = cell.text_frame.paragraphs[0].runs[0]
                    run.font.name = FONT
                    run.font.size = Pt(8.5)
                    run.font.color.rgb = self._rgb("ink")

        self._stat_bar(slide, stats)

    def _render_compliance(self, slide, sd: dict):
        """Compliance verdict badge + conditions/docs columns."""
        from pptx.util import Inches, Pt
        self._bg(slide)
        self._topbar(slide, sd.get("eyebrow", "Compliance Status"))

        verdict = (sd.get("verdict") or "CLEAR").upper()
        verdict_label = sd.get("verdict_label") or {
            "CLEAR": "CLEAR TO PROCEED",
            "CONDITIONS": "PROCEED WITH CONDITIONS",
            "BLOCKED": "BLOCKED",
        }.get(verdict, verdict)
        verdict_colors = {
            "CLEAR": ("DFF3F0", "0A7A71"),
            "CONDITIONS": ("FEF3C7", "92620E"),
            "BLOCKED": ("FEE2E2", "991B1B"),
        }
        bg_hex, fg_hex = verdict_colors.get(verdict, verdict_colors["CLEAR"])
        conditions = (sd.get("conditions") or [])[:6]
        docs = (sd.get("docs_required") or [])[:6]
        blocker = sd.get("blocker", "")
        stats = sd.get("stats") or []

        badge = slide.shapes.add_shape(1, Inches(PAD_X), Inches(BODY_TOP), Inches(3.6), Inches(0.42))
        badge.fill.solid()
        badge.fill.fore_color.rgb = self._hex_rgb(bg_hex)
        badge.line.fill.background()
        self._round_corners(badge, 20000)
        vb = slide.shapes.add_textbox(Inches(PAD_X + 0.14), Inches(BODY_TOP + 0.06),
                                      Inches(3.3), Inches(0.30))
        vp = vb.text_frame.paragraphs[0]
        vr = vp.add_run()
        vr.text = verdict_label
        vr.font.name = FONT
        vr.font.size = Pt(12)
        vr.font.bold = True
        vr.font.color.rgb = self._hex_rgb(fg_hex)

        y = BODY_TOP + 0.58
        if blocker:
            box = slide.shapes.add_shape(1, Inches(PAD_X), Inches(y), Inches(CONTENT_W), Inches(0.42))
            box.fill.solid()
            box.fill.fore_color.rgb = self._hex_rgb("FEE2E2")
            box.line.fill.background()
            self._round_corners(box, 8000)
            self._text(slide, PAD_X + 0.10, y + 0.06, CONTENT_W - 0.20, 0.30,
                       f"Blocker: {blocker}", 9.5, "ink", bold=True)
            y += 0.50

        col_w = (CONTENT_W - 0.24) / 2
        col_h = SLIDE_H - y - STAT_BAR_H - 0.10

        def _col(x, label, lst):
            self._text(slide, x, y, col_w, 0.20, label.upper(), 8, "gray", bold=True)
            cy = y + 0.26
            avail = col_h - 0.26
            # Real compliance/legal text routinely wraps to 2+ lines — size each item to its
            # actual wrapped height (not a fixed row height) so items never overlap, and drop
            # trailing items rather than overflow the slide if they don't all fit.
            kept, heights, more = self._fit_list_height(lst, col_w - 0.20, 8.5, avail)
            for it, h in zip(kept, heights):
                self._text(slide, x, cy, col_w, h - 0.02, f"•  {it}", 8.5, "ink")
                cy += h
            if more:
                self._text(slide, x, cy, col_w, 0.20, f"+{more} khác…", 8, "gray_lt", italic=True)

        if conditions:
            _col(PAD_X, "Conditions Before Launch", conditions)
        if docs:
            _col(PAD_X + col_w + 0.24, "Documents Required", docs)

        self._stat_bar(slide, stats)

    def _render_roi(self, slide, sd: dict):
        """ROI / Why Now — dark stat cards + reasons list."""
        from pptx.util import Inches, Pt
        self._bg(slide)
        self._topbar(slide, sd.get("eyebrow", "ROI / Why Now"))
        hl = sd.get("headline", {})
        self._text_mixed(slide, PAD_X, BODY_TOP, CONTENT_W, 0.42,
                         hl.get("plain", ""), hl.get("bold", ""), 18)

        stats = (sd.get("stats") or [])[:3]
        reasons = (sd.get("reasons") or [])[:5]
        color_map = {"orange": "orange", "teal": "teal", "purple": "purple", "gold": "gold"}

        body_top = BODY_TOP + 0.50
        body_h = SLIDE_H - body_top - STAT_BAR_H - 0.10
        stats_col_w = 2.4
        gap = 0.24

        # Stat cards (dark), stacked vertically on the left
        if stats:
            card_h = min((body_h - 0.12 * (len(stats) - 1)) / len(stats), 1.10)
            cy = body_top
            for s in stats:
                card = slide.shapes.add_shape(1, Inches(PAD_X), Inches(cy),
                                              Inches(stats_col_w), Inches(card_h))
                card.fill.solid()
                card.fill.fore_color.rgb = self._rgb("ink")
                card.line.fill.background()
                self._round_corners(card, 10000)
                color_key = color_map.get(s.get("color", "orange"), "orange")
                self._text(slide, PAD_X + 0.14, cy + 0.10, stats_col_w - 0.28, 0.42,
                           s.get("value", ""), 22, color_key, bold=True, auto_fit=True)
                self._text(slide, PAD_X + 0.14, cy + 0.10 + 0.44, stats_col_w - 0.28, 0.40,
                           s.get("label", ""), 8.5, "white")
                cy += card_h + 0.12

        # Reasons list on the right
        if reasons:
            rx = PAD_X + stats_col_w + gap
            rw = CONTENT_W - stats_col_w - gap
            reason_h = min((body_h - 0.10 * (len(reasons) - 1)) / len(reasons), 0.60)
            ry = body_top + max((body_h - (reason_h * len(reasons) + 0.10 * (len(reasons) - 1))) / 2, 0)
            for r in reasons:
                box = slide.shapes.add_shape(1, Inches(rx), Inches(ry), Inches(rw), Inches(reason_h))
                box.fill.solid()
                box.fill.fore_color.rgb = self._rgb("cream")
                box.line.color.rgb = self._rgb("line")
                box.line.width = Pt(0.5)
                self._round_corners(box, 10000)
                self._text(slide, rx + 0.14, ry, rw - 0.28, reason_h, r, 9, "ink", auto_fit=True)
                ry += reason_h + 0.10

    def _render_timeline(self, slide, sd: dict):
        """Phased implementation timeline (week -> label -> items)."""
        from pptx.util import Inches, Pt
        self._bg(slide)
        self._topbar(slide, sd.get("eyebrow", "Next Steps & Timeline"))
        hl = sd.get("headline", {})
        self._text_mixed(slide, PAD_X, BODY_TOP, CONTENT_W, 0.42,
                         hl.get("plain", ""), hl.get("bold", ""), 18)

        weeks = (sd.get("weeks") or [])[:6]
        stats = sd.get("stats") or []
        top = BODY_TOP + 0.46
        avail_h = BODY_H - 0.46
        items_w = CONTENT_W - 1.2

        # Each week's joined item text can wrap to 2+ lines — an equal fixed row_h overflows
        # into the divider/next row for any week with more than a short one-liner. Size every
        # row to its real content height, then scale all rows down proportionally if the deck
        # as a whole doesn't fit (still bounded, never overlapping).
        items_texts = [" · ".join(w.get("items") or []) for w in weeks]
        min_row_h = 0.30
        raw_heights = [
            max(0.26 + self._est_h(t, items_w, 8.5), min_row_h) for t in items_texts
        ]
        total_h = sum(raw_heights) or 1.0
        scale = min(avail_h / total_h, 1.0) if weeks else 1.0
        row_heights = [max(h * scale, min_row_h * scale) for h in raw_heights]

        y = top
        for w, items_text, row_h in zip(weeks, items_texts, row_heights):
            self._text(slide, PAD_X, y + 0.06, 1.1, row_h - 0.10,
                       w.get("week", ""), 8.5, "orange", bold=True)
            self._text(slide, PAD_X + 1.2, y + 0.02, items_w, 0.24,
                       w.get("label", ""), 10, "ink", bold=True)
            self._text(slide, PAD_X + 1.2, y + 0.26, items_w, row_h - 0.28,
                       items_text, 8.5, "gray", auto_fit=True)
            if row_h > 0.20:
                div = slide.shapes.add_shape(1, Inches(PAD_X), Inches(y + row_h - 0.02),
                                             Inches(CONTENT_W), Inches(0.008))
                div.fill.solid()
                div.fill.fore_color.rgb = self._rgb("line")
                div.line.fill.background()
            y += row_h

        self._stat_bar(slide, stats)

    def _render_checklist(self, slide, sd: dict):
        """Client decisions + tech-confirmation checklist."""
        from pptx.util import Inches, Pt
        self._bg(slide)
        self._topbar(slide, sd.get("eyebrow", "Key Decisions Required"))
        hl = sd.get("headline", {})
        self._text_mixed(slide, PAD_X, BODY_TOP, CONTENT_W, 0.42,
                         hl.get("plain", ""), hl.get("bold", ""), 18)

        decisions = (sd.get("decisions") or [])[:6]
        tech_items = (sd.get("tech_items") or [])[:6]
        stats = sd.get("stats") or []

        y0 = BODY_TOP + 0.50
        col_gap = 0.28
        two_col = bool(tech_items)
        col_w = (CONTENT_W - col_gap) / 2 if two_col else CONTENT_W

        col_h = SLIDE_H - y0 - STAT_BAR_H - 0.10
        avail = col_h - 0.26

        self._text(slide, PAD_X, y0, col_w, 0.20, "CLIENT DECISIONS", 8, "gray", bold=True)
        cy = y0 + 0.26
        # Decision text can run long (budget/approval sentences) — size each row to its real
        # wrapped height so rows never overlap; drop trailing items if they don't all fit.
        d_kept, d_heights, d_more = self._fit_list_height(
            decisions, col_w - 0.24, 9, avail, min_item_h=0.28)
        for d, h in zip(d_kept, d_heights):
            text = d.get("text", "") if isinstance(d, dict) else d
            priority = d.get("priority", "") if isinstance(d, dict) else ""
            box = slide.shapes.add_shape(1, Inches(PAD_X), Inches(cy), Inches(0.16), Inches(0.16))
            box.fill.solid()
            box.fill.fore_color.rgb = self._rgb("orange" if priority == "high" else "line")
            box.line.fill.background()
            self._text(slide, PAD_X + 0.24, cy - 0.03, col_w - 0.24, h - 0.02, str(text), 9, "ink")
            cy += h
        if d_more:
            self._text(slide, PAD_X, cy, col_w, 0.20, f"+{d_more} khác…", 8, "gray_lt", italic=True)

        if two_col:
            tx = PAD_X + col_w + col_gap
            self._text(slide, tx, y0, col_w, 0.20, "TECH CONFIRMATION REQUIRED", 8, "gray", bold=True)
            ty = y0 + 0.26
            t_kept, t_heights, t_more = self._fit_list_height(
                tech_items, col_w - 0.10, 9, avail, min_item_h=0.24)
            for t, h in zip(t_kept, t_heights):
                self._text(slide, tx, ty, col_w, h - 0.02, f"⚙  {t}", 9, "ink")
                ty += h
            if t_more:
                self._text(slide, tx, ty, col_w, 0.20, f"+{t_more} khác…", 8, "gray_lt", italic=True)

        self._stat_bar(slide, stats)

    def _render_value(self, slide, sd: dict):
        from pptx.util import Inches, Pt
        from pptx.util import Pt as _Pt
        self._bg(slide)
        hl = sd.get("headline", {})
        self._topbar(slide, sd.get("eyebrow", ""), sd.get("tier", ""))

        # Headline — 20pt; LEFT_W_V gives space for right stat col
        self._text_mixed(slide, PAD_X, BODY_TOP, LEFT_W_V, 1.10,
                         hl.get("plain", ""), hl.get("bold", ""), 20)

        # Lede — tighter gap after headline
        lede = sd.get("lede", "")
        lede_h = 0.40
        if lede:
            self._text(slide, PAD_X, BODY_TOP + 1.14, LEFT_W_V - 0.10, lede_h, lede, 9.5, "gray")

        # Feature cards — use full LEFT_W_V width
        cards = sd.get("cards") or []
        card_y = BODY_TOP + (1.58 if lede else 1.14)
        cards_available_h = BODY_H - (card_y - BODY_TOP)
        card_h = max(cards_available_h / max(len(cards), 1) - 0.06, 0.52)
        for c in cards[:4]:
            crd = slide.shapes.add_shape(1, Inches(PAD_X), Inches(card_y),
                                         Inches(LEFT_W_V), Inches(card_h))
            crd.fill.solid()
            crd.fill.fore_color.rgb = self._rgb("white")
            crd.line.color.rgb = self._rgb("line")
            crd.line.width = _Pt(0.5)
            self._round_corners(crd, 10000)

            # Icon box — rounder (HTML uses border-radius:8px on 32px = 25%)
            icon_box = slide.shapes.add_shape(1, Inches(PAD_X + 0.07), Inches(card_y + 0.06),
                                              Inches(0.28), Inches(0.28))
            icon_box.fill.solid()
            icon_box.fill.fore_color.rgb = self._rgb("cream")
            icon_box.line.fill.background()
            self._round_corners(icon_box, 20000)
            self._text(slide, PAD_X + 0.07, card_y + 0.05, 0.30, 0.28,
                       c.get("icon", ""), 11, "ink", align="CENTER", emoji_font=True)

            avail = card_h - 0.06
            title_h = min(avail * 0.55, 0.36)
            desc_h  = min(avail * 0.42, 0.28)
            self._text(slide, PAD_X + 0.43, card_y + 0.05, LEFT_W_V - 0.50, title_h,
                       c.get("title", ""), 9, "ink", bold=True, auto_fit=True)
            self._text(slide, PAD_X + 0.43, card_y + 0.05 + title_h + 0.02, LEFT_W_V - 0.50, desc_h,
                       c.get("desc", ""), 8.5, "gray", auto_fit=True)
            card_y += card_h + 0.05

        # Right stat column — mirrors HTML rs-card-primary / rs-card-secondary
        stats = sd.get("stats") or []
        if stats:
            # Primary card (orange bg)
            p_h = 1.36
            pcard = slide.shapes.add_shape(1, Inches(RC_X), Inches(BODY_TOP),
                                           Inches(RC_W), Inches(p_h))
            pcard.fill.solid()
            pcard.fill.fore_color.rgb = self._rgb("orange")
            pcard.line.fill.background()
            self._round_corners(pcard, 8000)
            s0 = stats[0]
            self._text(slide, RC_X + 0.14, BODY_TOP + 0.16, RC_W - 0.28, 0.54,
                       s0.get("v", ""), 24, "white", bold=True, auto_fit=True)
            self._text(slide, RC_X + 0.14, BODY_TOP + 0.74, RC_W - 0.28, 0.48,
                       s0.get("l", ""), 8.5, "white")

            # Secondary cards
            sec_y = BODY_TOP + p_h + 0.12
            for s in stats[1:3]:
                scard = slide.shapes.add_shape(1, Inches(RC_X), Inches(sec_y),
                                               Inches(RC_W), Inches(0.88))
                scard.fill.solid()
                scard.fill.fore_color.rgb = self._rgb("white")
                scard.line.color.rgb = self._rgb("line")
                scard.line.width = _Pt(0.5)
                self._round_corners(scard, 8000)
                self._text(slide, RC_X + 0.10, sec_y + 0.10, RC_W - 0.20, 0.36,
                           s.get("v", ""), 20, "ink", bold=True, auto_fit=True)
                self._text(slide, RC_X + 0.10, sec_y + 0.50, RC_W - 0.20, 0.26,
                           s.get("l", ""), 8, "gray_lt")
                sec_y += 0.88 + 0.12

    def _render_flow(self, slide, sd: dict):
        from pptx.util import Inches, Pt
        self._bg(slide)
        hl = sd.get("headline", {})
        self._topbar(slide, sd.get("eyebrow", ""))

        # Headline (full width)
        self._text_mixed(slide, PAD_X, BODY_TOP, CONTENT_W, 0.80,
                         hl.get("plain", ""), hl.get("bold", ""), 22)

        # ── Legend row ──
        lg_y = BODY_TOP + 0.80
        # Teal dot: Core
        d1 = slide.shapes.add_shape(1, Inches(PAD_X), Inches(lg_y + 0.04), Inches(0.10), Inches(0.10))
        d1.fill.solid(); d1.fill.fore_color.rgb = self._rgb("teal"); d1.line.fill.background()
        self._text(slide, PAD_X + 0.14, lg_y, 1.3, 0.18, "Core (có sẵn)", 8.5, "gray")
        # Orange dot: Custom
        d2 = slide.shapes.add_shape(1, Inches(PAD_X + 1.6), Inches(lg_y + 0.04), Inches(0.10), Inches(0.10))
        d2.fill.solid(); d2.fill.fore_color.rgb = self._rgb("orange"); d2.line.fill.background()
        self._text(slide, PAD_X + 1.74, lg_y, 1.5, 0.18, "Custom (mở rộng)", 8.5, "gray")
        # Role color explanation (right side)
        self._text(slide, PAD_X + 3.6, lg_y, 5.5, 0.18,
                   "Tím = Admin · CMS    Xanh = Customer · Mini App", 8.5, "gray")

        # ── Footer layout ──
        footer = sd.get("footer", "")
        stat_div_y = SLIDE_H - STAT_BAR_H - 0.02
        footer_h = 0.40
        footer_y = stat_div_y - footer_h - 0.08 if footer else stat_div_y

        # ── Steps ──
        steps = (sd.get("steps") or [])[:6]
        n = max(len(steps), 1)
        step_w = CONTENT_W / n
        flow_top = lg_y + 0.24
        icon_size = 0.60
        # Center icon row vertically between flow_top and footer_y
        label_h = 0.26 + 0.42  # label + desc below icon
        avail_v = footer_y - flow_top - 0.26 - icon_size - label_h
        icon_y = flow_top + 0.26 + max(avail_v / 2, 0)

        role_colors = {
            "customer": "teal", "admin": "purple",
            "staff": "orange2", "system": "gray_lt",
        }

        for i, st in enumerate(steps):
            cx = PAD_X + i * step_w + step_w / 2
            role = st.get("role", "customer")
            dot = st.get("dot", "core")
            rc = role_colors.get(role, "gray_lt")

            # Role pill — fully rounded (pill shape)
            pill_w, pill_h = min(step_w - 0.12, 0.90), 0.16
            px = cx - pill_w / 2
            pill = slide.shapes.add_shape(1, Inches(px), Inches(icon_y - 0.26),
                                           Inches(pill_w), Inches(pill_h))
            pill.fill.solid(); pill.fill.fore_color.rgb = self._rgb(rc)
            pill.line.fill.background()
            self._round_corners(pill, 50000)
            self._text(slide, px + 0.02, icon_y - 0.26, pill_w - 0.04, pill_h,
                       role.upper(), 6.5, "white", bold=True, align="CENTER")

            # Icon box — rounded (HTML: border-radius 16px on 68px ≈ 24%)
            ix = cx - icon_size / 2
            icon_box = slide.shapes.add_shape(1, Inches(ix), Inches(icon_y),
                                               Inches(icon_size), Inches(icon_size))
            icon_box.fill.solid(); icon_box.fill.fore_color.rgb = self._rgb("white")
            icon_box.line.color.rgb = self._rgb("line"); icon_box.line.width = Pt(0.5)
            self._round_corners(icon_box, 20000)
            self._text(slide, ix + 0.02, icon_y + 0.05, icon_size - 0.04, icon_size - 0.10,
                       st.get("icon", ""), 18, "ink", align="CENTER", emoji_font=True)

            # Core/custom dot — circle
            dot_color = "teal" if dot == "core" else "orange"
            dd = slide.shapes.add_shape(1, Inches(cx + icon_size / 2 - 0.10), Inches(icon_y - 0.04),
                                         Inches(0.11), Inches(0.11))
            dd.fill.solid(); dd.fill.fore_color.rgb = self._rgb(dot_color)
            dd.line.fill.background()
            self._round_corners(dd, 50000)

            # Label + desc centered below icon
            lbl_w = min(step_w - 0.10, 1.20)
            lx = cx - lbl_w / 2
            self._text(slide, lx, icon_y + icon_size + 0.10, lbl_w, 0.26,
                       st.get("label", ""), 9, "ink", bold=True, align="CENTER")
            self._text(slide, lx, icon_y + icon_size + 0.38, lbl_w, 0.42,
                       st.get("desc", ""), 7.5, "gray", align="CENTER")

            # Arrow between steps
            if i < n - 1:
                arr_x = cx + icon_size / 2 + 0.04
                arr_w = max(step_w - icon_size - 0.14, 0.05)
                arr = slide.shapes.add_shape(1, Inches(arr_x), Inches(icon_y + icon_size / 2 - 0.005),
                                              Inches(arr_w), Inches(0.01))
                arr.fill.solid(); arr.fill.fore_color.rgb = self._rgb("gray_lt")
                arr.line.fill.background()

        # ── Footer (if present) ──
        if footer:
            fdiv = slide.shapes.add_shape(1, Inches(PAD_X), Inches(footer_y),
                                           Inches(CONTENT_W), Inches(0.008))
            fdiv.fill.solid(); fdiv.fill.fore_color.rgb = self._rgb("line")
            fdiv.line.fill.background()
            tb = slide.shapes.add_textbox(Inches(PAD_X), Inches(footer_y + 0.04),
                                           Inches(CONTENT_W), Inches(footer_h))
            tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            r1 = p.add_run()
            r1.text = "Custom thêm: "
            r1.font.name = FONT; r1.font.size = Pt(8)
            r1.font.bold = True; r1.font.color.rgb = self._rgb("ink")
            r2 = p.add_run()
            r2.text = footer
            r2.font.name = FONT; r2.font.size = Pt(8)
            r2.font.bold = False; r2.font.color.rgb = self._rgb("gray")

        self._stat_bar(slide, sd.get("stats") or [], no_line=bool(footer))

    def _render_highlight(self, slide, sd: dict):
        """Exec-summary slide: headline + summary text + horizontal metric cards."""
        from pptx.util import Inches, Pt
        self._bg(slide)
        hl = sd.get("headline", {})
        self._topbar(slide, sd.get("eyebrow", ""))

        # Headline — 24pt; 1.00" tall holds 2 lines
        self._text_mixed(slide, PAD_X, BODY_TOP, CONTENT_W, 1.00,
                         hl.get("plain", ""), hl.get("bold", ""), 24)

        # Summary sentence below headline
        summary = sd.get("summary", "")
        if summary:
            self._text(slide, PAD_X, BODY_TOP + 1.04, CONTENT_W * 0.80, 0.44,
                       summary, 10, "gray")

        # Metric cards — horizontal row, fixed compact height
        metrics = (sd.get("metrics") or [])[:4]
        n = max(len(metrics), 1)
        card_y = BODY_TOP + 1.58
        # Cap at 2.20" — fills available space without looking sparse
        stat_top = SLIDE_H - STAT_BAR_H - 0.04
        card_h = min(stat_top - card_y - 0.10, 2.20)
        gap = 0.10
        card_w = (CONTENT_W - gap * (n - 1)) / n

        # Fixed content sizes; center value+label block vertically inside card (below accent bar)
        val_h = 0.62
        lbl_h = 0.32
        remaining = card_h - 0.06  # space below accent bar
        v_pad = max((remaining - val_h - 0.06 - lbl_h) / 2.0, 0.10)
        val_y = card_y + 0.06 + v_pad

        from pptx.util import Pt as _Pt
        for i, m in enumerate(metrics):
            cx = PAD_X + i * (card_w + gap)

            # Card background
            crd = slide.shapes.add_shape(1, Inches(cx), Inches(card_y),
                                         Inches(card_w), Inches(card_h))
            crd.fill.solid()
            crd.fill.fore_color.rgb = self._rgb("white")
            crd.line.color.rgb = self._rgb("line")
            crd.line.width = _Pt(0.5)
            self._round_corners(crd, 8000)

            # Colored accent bar at top of card (same radius so it sits flush inside card)
            _METRIC_COLORS = {"orange": "orange", "teal": "teal", "purple": "purple", "gold": "gold"}
            color_name = (m.get("color") or "orange")
            bar_color = self._rgb(_METRIC_COLORS.get(color_name, "orange"))
            bar = slide.shapes.add_shape(1, Inches(cx), Inches(card_y),
                                         Inches(card_w), Inches(0.06))
            bar.fill.solid()
            bar.fill.fore_color.rgb = bar_color
            bar.line.fill.background()
            self._round_corners(bar, 8000)

            # Value (big number) — vertically centered in card below accent bar
            self._text(slide, cx + 0.08, val_y, card_w - 0.16, val_h,
                       m.get("value", ""), 24, "ink", bold=True, align="CENTER", auto_fit=True)

            # Label — just below value
            self._text(slide, cx + 0.08, val_y + val_h + 0.06, card_w - 0.16, lbl_h,
                       m.get("label", ""), 8.5, "gray", align="CENTER", auto_fit=True)

        self._stat_bar(slide, sd.get("stats") or [])

    def _render_tier(self, slide, sd: dict):
        from pptx.util import Inches, Pt
        self._bg(slide)
        hl = sd.get("headline", {})
        self._topbar(slide, sd.get("eyebrow", ""))

        # Headline — 20pt keeps 2-line headlines within 0.82" so lede stays below
        self._text_mixed(slide, PAD_X, BODY_TOP, CONTENT_W, 0.82,
                         hl.get("plain", ""), hl.get("bold", ""), 20)

        lede = sd.get("lede", "")
        if lede:
            # Lede is safely below the headline block (headline ends at ~0.82")
            self._text(slide, PAD_X, BODY_TOP + 0.86, CONTENT_W - 1.0, 0.24, lede, 9.5, "gray")

        tiers = (sd.get("tiers") or [])[:4]
        n = max(len(tiers), 1)
        gap = 0.10  # gap between tier cards
        tier_w = (CONTENT_W - gap * (n - 1)) / n
        tier_top = BODY_TOP + (1.16 if lede else 0.86)
        tier_h = SLIDE_H - tier_top - STAT_BAR_H - 0.12

        for i, t in enumerate(tiers):
            x = PAD_X + i * (tier_w + gap)

            # Card bg
            crd = slide.shapes.add_shape(1, Inches(x + 0.04), Inches(tier_top),
                                         Inches(tier_w - 0.06), Inches(tier_h))
            crd.fill.solid()
            crd.fill.fore_color.rgb = self._rgb("white")
            crd.line.color.rgb = self._rgb("line")
            crd.line.width = Pt(0.5)
            self._round_corners(crd, 8000)

            # Top color bar (6px ≈ 0.06") — rounded top to match card
            bar_color = self._hex_rgb(t.get("barColor", "ECE6E1"))
            top_bar = slide.shapes.add_shape(1, Inches(x + 0.04), Inches(tier_top),
                                             Inches(tier_w - 0.06), Inches(0.06))
            top_bar.fill.solid()
            top_bar.fill.fore_color.rgb = bar_color
            top_bar.line.fill.background()
            self._round_corners(top_bar, 8000)

            # Tier name
            name_color = self._hex_rgb(t.get("nameColor", "1D1D1F"))
            name_tb = slide.shapes.add_textbox(Inches(x + 0.12), Inches(tier_top + 0.10),
                                               Inches(tier_w - 0.26), Inches(0.30))
            nf = name_tb.text_frame
            nf.word_wrap = True
            np_ = nf.paragraphs[0]
            nr = np_.add_run()
            nr.text = t.get("name", "")
            nr.font.name = FONT
            nr.font.size = Pt(9)
            nr.font.bold = True
            nr.font.color.rgb = name_color

            # Module
            self._text(slide, x + 0.12, tier_top + 0.38, tier_w - 0.26, 0.22,
                       t.get("module", ""), 7.5, "gray_lt", bold=True)

            # Price — split amount (large) + unit (small gray)
            price_str = t.get("price", "")
            p_parts = price_str.rsplit(" ", 1)
            amount_str = p_parts[0]
            unit_str = p_parts[1] if len(p_parts) > 1 else ""
            self._text(slide, x + 0.12, tier_top + 0.60, tier_w - 0.26, 0.38,
                       amount_str, 20, "ink", bold=True)
            if unit_str:
                self._text(slide, x + 0.12, tier_top + 0.96, tier_w - 0.26, 0.18,
                           unit_str, 8, "gray")

            # Period
            self._text(slide, x + 0.12, tier_top + 1.16, tier_w - 0.26, 0.22,
                       t.get("period", ""), 8, "gray_lt")

            # Divider
            d = slide.shapes.add_shape(1, Inches(x + 0.12), Inches(tier_top + 1.38),
                                       Inches(tier_w - 0.26), Inches(0.01))
            d.fill.solid()
            d.fill.fore_color.rgb = self._rgb("line")
            d.line.fill.background()

            # Checks — sized to real wrapped height so a longer feature line never overlaps
            # the next check or the deploy text below.
            deploy = t.get("deploy", "")
            deploy_y = tier_top + tier_h - 0.34
            checks_start = tier_top + 1.46
            available_h = deploy_y - checks_start - 0.10
            kept, heights, more = self._fit_list_height(
                (t.get("checks") or [])[:4], tier_w - 0.38 - 0.20, 8.5, available_h,
                min_item_h=0.24)
            cy = checks_start
            for ck, h in zip(kept, heights):
                self._text(slide, x + 0.24, cy, tier_w - 0.38, h - 0.02, "✓  " + ck, 8.5, "ink")
                cy += h

            # Deploy
            if deploy:
                self._text(slide, x + 0.12, deploy_y, tier_w - 0.26, 0.30,
                           f"Triển khai: {deploy}", 8, "gray")

        self._stat_bar(slide, sd.get("stats") or [])

    def _render_screen(self, slide, sd: dict):
        """Phone mockup slide — up to 3 phones side by side."""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        self._bg(slide)
        hl = sd.get("headline", {})
        self._topbar(slide, sd.get("eyebrow", ""))

        self._text_mixed(slide, PAD_X, BODY_TOP, CONTENT_W, 0.62,
                         hl.get("plain", ""), hl.get("bold", ""), 18)

        lede = sd.get("lede", "")
        title_h = 0.62 + (0.26 if lede else 0)
        if lede:
            self._text(slide, PAD_X, BODY_TOP + 0.64, CONTENT_W, 0.26, lede, 9, "gray")

        screens = (sd.get("screens") or [])[:3]
        n = max(len(screens), 1)

        phone_top = BODY_TOP + title_h + 0.10
        # 0.32" padding reserves space for phone label (0.18") below phone + gap to stat divider
        phone_area_h = SLIDE_H - phone_top - STAT_BAR_H - 0.32
        phone_h = min(phone_area_h, 3.10)
        phone_w = phone_h * 0.48

        total_w = n * phone_w
        gap = (CONTENT_W - total_w) / (n + 1)

        for i, sc in enumerate(screens):
            px = PAD_X + gap * (i + 1) + phone_w * i
            py = phone_top

            # Phone body
            body = slide.shapes.add_shape(5, Inches(px), Inches(py),
                                           Inches(phone_w), Inches(phone_h))
            body.fill.solid()
            body.fill.fore_color.rgb = RGBColor(26, 26, 46)
            body.line.fill.background()

            # Screen area
            margin = phone_w * 0.07
            scr_x = px + margin
            scr_y = py + phone_h * 0.055
            scr_w = phone_w - margin * 2
            scr_h = phone_h * 0.875

            scr = slide.shapes.add_shape(1, Inches(scr_x), Inches(scr_y),
                                         Inches(scr_w), Inches(scr_h))
            scr.fill.solid()
            scr.fill.fore_color.rgb = RGBColor(255, 255, 255)
            scr.line.fill.background()

            # Notch
            notch_w = scr_w * 0.25
            notch = slide.shapes.add_shape(
                1, Inches(scr_x + (scr_w - notch_w) / 2), Inches(scr_y),
                Inches(notch_w), Inches(0.06))
            notch.fill.solid()
            notch.fill.fore_color.rgb = RGBColor(26, 26, 46)
            notch.line.fill.background()

            # App bar
            bar_h = 0.20
            bar = slide.shapes.add_shape(1, Inches(scr_x), Inches(scr_y + 0.06),
                                         Inches(scr_w), Inches(bar_h))
            bar.fill.solid()
            bar.fill.fore_color.rgb = self._rgb("orange")
            bar.line.fill.background()
            self._text(slide, scr_x, scr_y + 0.06, scr_w, bar_h,
                       self._safe_text(sc.get("app_name", "")), 6, "white",
                       bold=True, align="CENTER")

            # Content items
            cy = scr_y + 0.06 + bar_h + 0.04
            cx_i = scr_x + 0.04
            iw = scr_w - 0.08
            scr_bottom = scr_y + scr_h - 0.06

            for item in (sc.get("items") or [])[:6]:
                if cy + 0.18 > scr_bottom:
                    break
                kind = item.get("kind", "row")

                if kind == "banner":
                    h = min(0.32, scr_bottom - cy - 0.04)
                    bshape = slide.shapes.add_shape(1, Inches(cx_i), Inches(cy),
                                                    Inches(iw), Inches(h))
                    bshape.fill.solid()
                    bshape.fill.fore_color.rgb = self._rgb("orange")
                    bshape.line.fill.background()
                    self._round_corners(bshape, 12000)
                    txt = self._safe_text(f"{item.get('emoji','')} {item.get('text','')}")
                    self._text(slide, cx_i + 0.02, cy + 0.04, iw - 0.04, h - 0.06,
                               txt, 5.5, "white", bold=True, align="CENTER")
                    cy += h + 0.04

                elif kind == "row":
                    h = 0.24
                    rshape = slide.shapes.add_shape(1, Inches(cx_i), Inches(cy),
                                                    Inches(iw), Inches(h))
                    rshape.fill.solid()
                    rshape.fill.fore_color.rgb = RGBColor(248, 248, 248)
                    rshape.line.color.rgb = RGBColor(236, 230, 225)
                    rshape.line.width = Pt(0.5)
                    self._round_corners(rshape, 12000)
                    self._text(slide, cx_i + 0.02, cy + 0.03, 0.18, 0.20,
                               self._safe_text(item.get("emoji", "")), 8, "ink",
                               align="CENTER", emoji_font=True)
                    self._text(slide, cx_i + 0.22, cy + 0.04, iw - 0.28, 0.18,
                               self._safe_text(item.get("title", "")), 5.5, "ink")
                    cy += h + 0.03

                elif kind == "cta":
                    h = 0.22
                    cshape = slide.shapes.add_shape(1, Inches(cx_i + 0.04), Inches(cy),
                                                    Inches(iw - 0.08), Inches(h))
                    cshape.fill.solid()
                    cshape.fill.fore_color.rgb = self._rgb("orange")
                    cshape.line.fill.background()
                    self._round_corners(cshape, 50000)
                    self._text(slide, cx_i + 0.06, cy + 0.04, iw - 0.12, h - 0.06,
                               self._safe_text(item.get("text", "")), 6, "white",
                               bold=True, align="CENTER")
                    cy += h + 0.04

                elif kind == "zns":
                    h = 0.34
                    if cy + h > scr_bottom:
                        break
                    zshape = slide.shapes.add_shape(1, Inches(cx_i), Inches(cy),
                                                    Inches(iw), Inches(h))
                    zshape.fill.solid()
                    zshape.fill.fore_color.rgb = RGBColor(255, 255, 255)
                    zshape.line.color.rgb = self._rgb("orange")
                    zshape.line.width = Pt(1.0)
                    self._round_corners(zshape, 12000)
                    self._text(slide, cx_i + 0.04, cy + 0.03, iw - 0.08, 0.14,
                               self._safe_text(item.get("title", "")), 5.5, "orange",
                               bold=True)
                    self._text(slide, cx_i + 0.04, cy + 0.17, iw - 0.08, 0.14,
                               self._safe_text(item.get("text", "")), 5, "gray")
                    cy += h + 0.04

                elif kind == "points":
                    h = 0.22
                    pshape = slide.shapes.add_shape(1, Inches(cx_i), Inches(cy),
                                                    Inches(iw), Inches(h))
                    pshape.fill.solid()
                    pshape.fill.fore_color.rgb = self._rgb("teal")
                    pshape.line.fill.background()
                    self._round_corners(pshape, 12000)
                    txt = self._safe_text(
                        f"{item.get('emoji','⭐')} {item.get('value','')} {item.get('text','')}")
                    self._text(slide, cx_i + 0.02, cy + 0.04, iw - 0.04, h - 0.06,
                               txt, 6, "white", bold=True, align="CENTER")
                    cy += h + 0.03

            # Screen label below phone
            self._text(slide, px, py + phone_h + 0.04, phone_w, 0.18,
                       self._safe_text(sc.get("label", f"Screen {i + 1}")),
                       7, "gray_lt", align="CENTER")

        self._stat_bar(slide, sd.get("stats") or [])


def create_adtimabox_pptx_generator() -> AdtimaBoxPPTXGenerator:
    return AdtimaBoxPPTXGenerator()
