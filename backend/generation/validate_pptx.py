"""
Pre-export validator for CorporatePPTXGenerator's output.

Ported from agents/adtimabox-proposal-builder/reference/validate_deck.py (a
pptxgenjs-based reference deck's validator) — generalized to check structural
invariants of ANY deck this generator produces, not one fixed sample: this
version parses real slide content instead of comparing against hardcoded
reference dollar amounts and topic strings, so it runs against live output.

Geometry constants are imported from pptx_corporate.py rather than copied, so
this validator can't silently drift from the renderer it checks.

Usage:
    python -m generation.validate_pptx path/to/deck.pptx
or:
    from generation.validate_pptx import validate
    findings = validate("path/to/deck.pptx")   # {} means all checks passed
"""

from __future__ import annotations

import re
import sys
import hashlib
import collections

from pptx import Presentation
from pptx.util import Emu

from generation.pptx_corporate import (
    SLIDE_W, SLIDE_H, PAD_X, CONTENT_W, BODY_TOP, _STATIC_INTRO, _AGENDA_ITEMS,
)

R = PAD_X + CONTENT_W
CONTENT_BOT = 6.55
BOT_LIMIT = 7.05
FONT_FLOOR = 11.0
WF = 0.58  # must match pptx_corporate.py's _est_lines — see that module's docstring
EMOJI_RE = re.compile("[\U0001F300-\U0001FAFF☀-➿⬀-⯿️]")
PLACEHOLDER_TOKENS = ("[", "TBD", "N/A", "undefined", "null", "lorem", "{{")


def _is_fixed_slide(slide) -> bool:
    """A full-bleed picture with no text shapes — the static intro/cover pattern."""
    pics = [sh for sh in slide.shapes if "PICTURE" in str(sh.shape_type)]
    texts = [sh for sh in slide.shapes if sh.has_text_frame and sh.text_frame.text.strip()]
    return len(pics) == 1 and not texts


def validate(path: str) -> dict[int, list[str]]:
    """Run every check against the deck at `path`. Returns {check_num: [failure, ...]}
    — empty dict means everything passed."""
    prs = Presentation(path)
    slides = list(prs.slides)
    fails: dict[int, list[str]] = collections.defaultdict(list)
    texts_over_40: list[tuple[int, str]] = []

    for i, slide in enumerate(slides, 1):
        for sh in slide.shapes:
            left, top = Emu(sh.left).inches, Emu(sh.top).inches
            w = Emu(sh.width).inches
            h = sum(Emu(r.height).inches for r in sh.table.rows) if sh.has_table else Emu(sh.height).inches

            # 1. Canvas bounds (full-bleed images at 0,0,SLIDE_W,SLIDE_H are exempt)
            is_full_bleed = abs(left) < 0.01 and abs(w - SLIDE_W) < 0.01
            if (left + w > R + 0.02 or top + h > BOT_LIMIT + 0.02) and not is_full_bleed:
                fails[1].append(f"s{i} {sh.shape_type} right={left+w:.2f} bottom={top+h:.2f}")

            if not sh.has_text_frame:
                continue
            tf = sh.text_frame
            txt = tf.text

            # 2. autofit / margins / wrap — see pptx_corporate.py's _text() for why
            if tf.auto_size is not None:
                fails[2].append(f"s{i} autofit={tf.auto_size}")
            if txt.strip():
                if any(m not in (0, None) for m in
                       (tf.margin_top, tf.margin_bottom, tf.margin_left, tf.margin_right)):
                    fails[2].append(f"s{i} margin!=0 {txt[:30]!r}")
                if tf.word_wrap is False:
                    fails[2].append(f"s{i} wrap=False {txt[:30]!r}")

            maxpt = 0.0
            for para in tf.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        pt = run.font.size.pt
                        maxpt = max(maxpt, pt)
                        # 3. font floor
                        if pt < FONT_FLOOR:
                            fails[3].append(f"s{i} {pt}pt {run.text[:24]!r}")
                    # 10. no emoji
                    if EMOJI_RE.search(run.text or ""):
                        fails[10].append(f"s{i} {run.text[:24]!r}")

            # 4. estimated text height <= box height
            if txt.strip() and maxpt:
                cpl = max(1, int(w * 72 / (maxpt * WF)))
                lines = max(1, -(-len(txt) // cpl))
                need = (lines * maxpt * 1.2 + 6) / 72
                if need > h + 0.02:
                    fails[4].append(f"s{i} needs {need:.2f}in, box is {h:.2f}in {txt[:36]!r}")

            # 6. placeholder text
            if txt.strip() and any(tok in txt for tok in PLACEHOLDER_TOKENS):
                fails[6].append(f"s{i} placeholder-like text {txt[:36]!r}")

            # 11. double space / orphan trailing separator
            if txt.strip():
                if "  " in txt or txt.strip().endswith(("·", "-")):
                    fails[11].append(f"s{i} {txt[:40]!r}")
                if len(txt.strip()) > 40:
                    texts_over_40.append((i, txt.strip()))

    # 5. duplicate long strings across slides
    seen: dict[str, int] = {}
    for i, t in texts_over_40:
        key = hashlib.md5(t.encode()).hexdigest()
        if key in seen and seen[key] != i:
            fails[5].append(f"s{seen[key]} == s{i}: {t[:48]!r}")
        seen.setdefault(key, i)

    # 7. agenda lists exactly len(_AGENDA_ITEMS) sections, ascending, starting
    # right after the static intro block, none pointing past the deck's end.
    # _AGENDA_ITEMS[0] is the static intro block itself ("Slides 03 to 07"), so
    # its range must start at slide 3; every following row's numbers must be
    # ascending and land inside the deck.
    agenda_slide = slides[1] if len(slides) > 1 else None
    if agenda_slide is not None:
        agenda_texts = [sh.text_frame.text for sh in agenda_slide.shapes if sh.has_text_frame]
        ranges = [t for t in agenda_texts if t.startswith("Slide")]
        nums = [int(x) for t in ranges for x in re.findall(r"\d+", t)]
        if len(ranges) != len(_AGENDA_ITEMS):
            fails[7].append(f"agenda lists {len(ranges)} section ranges, expected {len(_AGENDA_ITEMS)}")
        if nums and nums[0] != 3:
            fails[7].append(f"static-intro range starts at slide {nums[0]}, expected 3")
        if nums and max(nums) > len(slides):
            fails[7].append(f"agenda references slide {max(nums)} but deck has {len(slides)}")
        if nums != sorted(nums):
            fails[7].append("agenda ranges not ascending")

    # 8. quotation money maths — find this generator's own "Subtotal excl. VAT" /
    # "VAT 8%" / "TOTAL" captions and verify subtotal*8%==vat, subtotal+vat==total
    def _money(s: str):
        t = s.strip().replace(".", "").replace(",", "")
        neg = t.startswith("-")
        t = t.lstrip("-")
        return (-1 if neg else 1) * int(t) if t.isdigit() else None

    for i, slide in enumerate(slides, 1):
        labels = {sh.text_frame.text.strip(): sh for sh in slide.shapes if sh.has_text_frame}
        if "Subtotal excl. VAT" not in labels or "VAT 8%" not in labels:
            continue
        # values sit in sibling shapes at the same row (same top, to the right)
        def _value_right_of(label_shape):
            ly = Emu(label_shape.top).inches
            candidates = [
                sh for sh in slide.shapes
                if sh.has_text_frame and abs(Emu(sh.top).inches - ly) < 0.03
                and Emu(sh.left).inches > Emu(label_shape.left).inches
            ]
            return _money(candidates[0].text_frame.text) if candidates else None

        subtotal = _value_right_of(labels["Subtotal excl. VAT"])
        vat = _value_right_of(labels["VAT 8%"])
        total = _value_right_of(labels["TOTAL"]) if "TOTAL" in labels else None
        if subtotal is not None and vat is not None:
            if round(subtotal * 0.08) != vat:
                fails[8].append(f"s{i} VAT mismatch: subtotal {subtotal} * 8% != {vat}")
        if subtotal is not None and vat is not None and total is not None:
            if subtotal + vat != total:
                fails[8].append(f"s{i} total mismatch: {subtotal}+{vat} != {total}")

    # 13. content density 60-92% (cover, closing, fixed intro slides exempt)
    for i, slide in enumerate(slides, 1):
        if i in (1, len(slides)) or _is_fixed_slide(slide):
            continue
        bottoms = []
        for sh in slide.shapes:
            t = Emu(sh.top).inches
            hh = sum(Emu(r.height).inches for r in sh.table.rows) if sh.has_table else Emu(sh.height).inches
            if t >= BODY_TOP - 0.05 and t + hh <= CONTENT_BOT + 0.05:
                bottoms.append(t + hh)
        if not bottoms:
            fails[13].append(f"s{i} no content in the content band")
            continue
        fill = (max(bottoms) - BODY_TOP) / (CONTENT_BOT - BODY_TOP)
        if fill < 0.60 or fill > 0.92:
            fails[13].append(f"s{i} content fills {fill*100:.0f}% (want 60-92%)")

    # 14. no two text boxes overlap by more than 0.05in on both axes
    for i, slide in enumerate(slides, 1):
        boxes = []
        for sh in slide.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                x1, y1 = Emu(sh.left).inches, Emu(sh.top).inches
                boxes.append((sh.text_frame.text, x1, y1, x1 + Emu(sh.width).inches, y1 + Emu(sh.height).inches))
        for a in range(len(boxes)):
            for b in range(a + 1, len(boxes)):
                t1, x1, y1, x2, y2 = boxes[a]
                t2, X1, Y1, X2, Y2 = boxes[b]
                ox, oy = min(x2, X2) - max(x1, X1), min(y2, Y2) - max(y1, Y1)
                if ox > 0.05 and oy > 0.05:
                    fails[14].append(f"s{i} {t1[:24]!r} overlaps {t2[:24]!r} ({ox:.2f}x{oy:.2f}in)")

    # 15. fixed collateral block: len(_STATIC_INTRO) full-bleed slides right
    # after cover + agenda, in order
    fixed_idx = [i for i, sl in enumerate(slides, 1) if _is_fixed_slide(sl)]
    want = list(range(3, 3 + len(_STATIC_INTRO)))
    if fixed_idx[:len(want)] != want:
        fails[15].append(f"static intro block at slides {fixed_idx}, expected {want}")

    return dict(fails)


_NAMES = {
    1: "Canvas bounds", 2: "Autofit/margin/wrap", 3: f"Font floor {FONT_FLOOR}pt",
    4: "Text fits its box", 5: "No duplicate long strings", 6: "No placeholder text",
    7: "Agenda matches deck", 8: "Quotation money maths", 10: "No emoji",
    11: "No double space / orphan separator", 13: "Content density 60-92%",
    14: "No overlapping text boxes", 15: "Static collateral block present",
}


def main() -> int:
    if len(sys.argv) < 2:
        print("usage: python -m generation.validate_pptx path/to/deck.pptx")
        return 2
    path = sys.argv[1]
    fails = validate(path)
    prs = Presentation(path)
    print(f"Deck: {path}  |  {len(prs.slides.__iter__.__self__._sldIdLst)} slides\n")
    ok = True
    for num, name in _NAMES.items():
        v = fails.get(num, [])
        print(f"  {'PASS' if not v else 'FAIL'}  {num:>2}. {name}")
        for line in v[:6]:
            print(f"          - {line}")
        if v:
            ok = False
    print("\nRESULT:", "ALL CHECKS PASSED" if ok else "FIX REQUIRED")
    return 0 if ok else 1


if __name__ == "__main__":
    sys.exit(main())
