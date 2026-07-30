"""
Adtima corporate-branded PPTX generator.

Follows the visual skeleton of generation/sample-output.pptx: picture-backed
chrome (assets/zbh-*.jpg extracted from that sample) and the 5 static Zalo/Adtima
platform intro slides embedded verbatim. The client-specific section scheme
covers all 7 sections of the source proposal document (Executive Summary and
Next Steps were added beyond the sample's own 5-section agenda so no part of
the proposal goes unexploited — compliance still folds into Solution, matching
the sample). This is independent from html_deck.py's 14-type/7-section schema,
which the HTML deck keeps using unchanged — only the downloadable PPTX uses
this template.

Section/page numbering is computed here from the final slide order, never taken
from the extraction — slide counts vary (e.g. whether an alternative quotation
tier exists), so trusting an LLM to count them produces wrong numbers.
"""

from __future__ import annotations

import math
import os
from typing import Any

_HERE = os.path.dirname(os.path.abspath(__file__))
_ASSETS_DIR = os.path.join(_HERE, "assets")

SLIDE_W = 13.333
SLIDE_H = 7.5
FONT = "Be Vietnam Pro"

PAD_X = 0.85
CONTENT_W = SLIDE_W - PAD_X * 2   # 11.633"
BODY_TOP = 1.95                  # below title, matches the sample's card-grid top
TITLE_Y = 0.82
TITLE_H = 0.60

_C = {
    "blue": "0067F7",
    "cyan": "00A8F0",
    "red": "E70207",
    "navy": "012B81",
    "ink": "111318",
    "gray": "4A4F5A",
    "border": "DCE4EE",
    "white": "FFFFFF",
    "subtitle": "C9D8F5",
    "footnote": "9AA6B8",
}

# The 5 static Zalo/Adtima platform slides always sit at "Slides 03 to 07" — fixed
# marketing collateral, not derived from any client brief.
_STATIC_INTRO = [
    "zbh-s01-about-zalo.jpg",
    "zbh-s02-lifecycle-role.jpg",
    "zbh-s03-usp.jpg",
    "zbh-s04-public-to-private.jpg",
    "zbh-s05-importance-private.jpg",
]

# Canonical order for dynamic content slides + which section/eyebrow they belong to.
# One slide per type is kept (first occurrence) — extraction is told to skip any
# type it has no source content for, so a shorter proposal just omits that slide
# rather than fabricating one. All 7 sections of the original proposal document
# are represented here (Executive Summary and Next Steps were missing from the
# sample's own 5-section agenda — added as Section 2 and Section 7 so no part of
# the source proposal goes unexploited).
_CONTENT_ORDER = [
    ("executive_summary", 2, "EXECUTIVE SUMMARY", "Executive summary"),
    ("client_requirements", 3, "CLIENT REQUIREMENTS", "What the client needs"),
    ("solution_package", 4, "SOLUTION", "Package and delivery items"),
    ("user_journey", 4, "SOLUTION", "The user journey"),
    ("solution_flowchart", 4, "SOLUTION", "Solution flow diagram"),
    ("touchpoints_table", 4, "SOLUTION", "Automated touchpoints"),
    ("compliance", 4, "SOLUTION", "Compliance and launch conditions"),
    ("quotation", 5, "QUOTATION", "Quotation, recommended option"),
    ("quotation_alternative", 5, "QUOTATION", "Alternative option"),
    ("case_study", 6, "CASE STUDY", "Comparable campaigns we delivered"),
    ("next_steps", 7, "NEXT STEPS", "Next steps and timeline"),
]

_AGENDA_ITEMS = [
    (1, "About Zalo & Zalo Brand Hub", "Platform scale, the user lifecycle, our USP and why private traffic matters"),
    (2, "Executive summary", "The headline ask, key numbers and why now"),
    (3, "Client requirements", "Current state, core pain, desired outcome and the gap to close"),
    (4, "Solution", "Package, user journey, solution flow, messaging and launch conditions"),
    (5, "Quotation", "Grouped line items, totals and the alternative option"),
    (6, "Case study", "Comparable campaigns Adtima has delivered"),
    (7, "Next steps", "Implementation timeline, decisions and tech confirmation"),
]


class CorporatePPTXGenerator:
    """Generates the Adtima-corporate-branded PPTX from extracted slide JSON."""

    async def generate(self, proposal_text: str, brief: dict, output_path: str) -> dict:
        """Standalone entry point (extraction + build). The live pipeline calls
        the extraction function directly and passes slides_data to _build_pptx,
        but this stays for ad-hoc/manual use."""
        try:
            from pptx import Presentation  # noqa
        except ImportError:
            return {"status": "error", "error": "python-pptx not installed"}

        from generation.pptx_corporate_extract import extract_pptx_slides
        try:
            slides_data = await extract_pptx_slides(proposal_text, brief)
        except Exception as e:
            print(f"[CorporatePPTX] Extraction failed ({e}), building with empty content slides")
            slides_data = []

        try:
            prs = self._build_pptx(slides_data)
            os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
            prs.save(output_path)
            return {"status": "success", "file_path": output_path}
        except Exception as e:
            return {"status": "error", "error": str(e)}

    # ── Build ────────────────────────────────────────────────────────────

    def _build_pptx(self, slides_data: list[dict]):
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        prs.slide_width = Inches(SLIDE_W)
        prs.slide_height = Inches(SLIDE_H)
        blank = prs.slide_layouts[6]

        by_type: dict[str, dict] = {}
        for sd in slides_data:
            t = sd.get("type")
            if t and t not in by_type:
                by_type[t] = sd

        cover = by_type.get("cover", {})
        closing = by_type.get("closing", {})

        content_entries = [
            (ctype, section_num, section_label, fallback_title)
            for (ctype, section_num, section_label, fallback_title) in _CONTENT_ORDER
            if ctype in by_type
        ]

        # Page numbers: 1 for cover, 2 for agenda, 3-7 for the static intro (no
        # overlay, matching the sample), then sequential for every content slide —
        # i.e. always (index within final slide list) + 1. Cover/closing show none.
        # Bug fixed: this used to start at 2 and increment per content slide,
        # producing "03, 04, 05..." regardless of how many static intro slides
        # came before them — off by len(_STATIC_INTRO) from the real slide
        # position, and from the ranges _render_agenda computes for the same
        # slides (which does account for the static intro), so the agenda said
        # "Slide 10" while that same slide's own footer showed "05".
        page = 2 + len(_STATIC_INTRO)  # cover=1, agenda=2, static intro=3..(2+n_static)
        total_pages = 2 + len(_STATIC_INTRO) + len(content_entries) + 1  # + closing

        dispatch = {
            "executive_summary": self._render_executive_summary,
            "client_requirements": self._render_client_requirements,
            "solution_package": self._render_solution_package,
            "user_journey": self._render_user_journey,
            "solution_flowchart": self._render_solution_flowchart,
            "touchpoints_table": self._render_touchpoints_table,
            "compliance": self._render_compliance,
            "quotation": self._render_quotation,
            "quotation_alternative": self._render_quotation_alternative,
            "case_study": self._render_case_study,
            "next_steps": self._render_next_steps,
        }

        # 1. Cover
        slide = prs.slides.add_slide(blank)
        self._render_cover(slide, cover)

        # 2. Agenda — slide ranges computed from the actual content_entries list
        slide = prs.slides.add_slide(blank)
        self._render_agenda(slide, content_entries, len(_STATIC_INTRO))

        # 3-7. Static Zalo/Adtima intro slides, verbatim
        for filename in _STATIC_INTRO:
            slide = prs.slides.add_slide(blank)
            self._bg_picture(slide, filename)

        # 8+. Client-specific content slides
        for ctype, section_num, section_label, fallback_title in content_entries:
            page += 1
            slide = prs.slides.add_slide(blank)
            sd = by_type[ctype]
            self._bg_picture(slide, "zbh-frame.jpg")
            self._chrome(slide, section_num, section_label, page)
            self._title(slide, sd.get("title") or fallback_title)
            dispatch[ctype](slide, sd)

        # Last. Closing
        slide = prs.slides.add_slide(blank)
        self._render_closing(slide, closing)

        return prs

    # ── Low-level helpers ────────────────────────────────────────────────

    @staticmethod
    def _rgb(hex_str: str):
        from pptx.dml.color import RGBColor
        h = hex_str.lstrip("#").upper()
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    @staticmethod
    def _safe_text(text: str) -> str:
        """Strip characters that break OOXML: lone surrogates, variation selectors,
        and C0/C1 controls."""
        if not text:
            return ""
        result = []
        for ch in str(text):
            cp = ord(ch)
            if 0xD800 <= cp <= 0xDFFF:
                continue
            if 0xFE00 <= cp <= 0xFE0F or 0xE0100 <= cp <= 0xE01EF:
                continue
            if cp < 0x20 and ch not in "\t\n\r":
                continue
            result.append(ch)
        return "".join(result)

    @staticmethod
    def _est_lines(text: str, width_in: float, font_pt: float, bold: bool = False) -> int:
        """Estimate wrapped line count — python-pptx text boxes don't reflow their
        container, so callers must size boxes to real content height. Bold glyphs
        run wider, so a bold label needs a smaller chars-per-line estimate or it
        under-counts wraps and the next element crowds it (measured against
        rendered PPTX, not just this formula)."""
        if not text:
            return 0
        avg_char_w_in = font_pt * (0.60 if bold else 0.52) / 72.0
        chars_per_line = max(int(width_in / avg_char_w_in), 1)
        lines = 0
        for para in str(text).split("\n"):
            lines += max(1, math.ceil(len(para) / chars_per_line))
        return lines

    def _est_h(self, text: str, width_in: float, font_pt: float, line_factor: float = 1.22,
               bold: bool = False) -> float:
        return self._est_lines(text, width_in, font_pt, bold) * font_pt * line_factor / 72.0

    def _est_w(self, text: str, font_pt: float, bold: bool = True, pad_in: float = 0.36) -> float:
        """Estimate a pill/box width (inches) that fits `text` on one line —
        used for the section eyebrow pill, which the sample sizes to its label."""
        avg_char_w_in = font_pt * (0.60 if bold else 0.52) / 72.0
        return len(text or "") * avg_char_w_in + pad_in

    def _fit_list_height(self, items: list, width_in: float, font_pt: float, max_h_in: float,
                          item_gap: float = 0.06, min_item_h: float = 0.22,
                          line_factor: float = 1.22):
        """Greedily keep as many `items` as fit within `max_h_in`, sizing each to
        its real wrapped height. Returns (kept, heights, dropped_count)."""
        heights: list[float] = []
        kept: list = []
        total = 0.0
        for it in items:
            text = it.get("text", it) if isinstance(it, dict) else it
            h = max(self._est_h(text, width_in, font_pt, line_factor) + item_gap, min_item_h)
            if total + h > max_h_in and kept:
                break
            kept.append(it)
            heights.append(h)
            total += h
        return kept, heights, len(items) - len(kept)

    @staticmethod
    def _round_corners(shape, adj: int = 1744):
        """adj — corner radius as a fraction of the shorter edge x 100000. The
        sample uses ~1744 (1.7%) for cards, far less rounded than a typical pill."""
        from lxml import etree
        p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
        a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        sp = shape._element
        spPr = sp.find(f"{{{p_ns}}}spPr")
        if spPr is None:
            return
        prstGeom = spPr.find(f"{{{a_ns}}}prstGeom")
        if prstGeom is None:
            return
        prstGeom.set("prst", "roundRect")
        for old in prstGeom.findall(f"{{{a_ns}}}avLst"):
            prstGeom.remove(old)
        avLst = etree.SubElement(prstGeom, f"{{{a_ns}}}avLst")
        gd = etree.SubElement(avLst, f"{{{a_ns}}}gd")
        gd.set("name", "adj")
        gd.set("fmla", f"val {adj}")

    def _text(self, slide, left, top, width, height, text, size, color_hex,
              bold=False, italic=False, align="LEFT", auto_fit=False, spc=None):
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        aligns = {"LEFT": PP_ALIGN.LEFT, "CENTER": PP_ALIGN.CENTER, "RIGHT": PP_ALIGN.RIGHT}
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = tb.text_frame
        tf.word_wrap = True
        if auto_fit:
            from pptx.enum.text import MSO_AUTO_SIZE
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        p = tf.paragraphs[0]
        p.alignment = aligns.get(align, PP_ALIGN.LEFT)
        run = p.add_run()
        run.text = self._safe_text(text or "")
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = self._rgb(color_hex)
        run.font.bold = bold
        run.font.italic = italic
        if spc is not None:
            run.font._rPr.set("spc", str(spc))
        return tb

    def _rect(self, slide, left, top, width, height, fill_hex=None, line_hex=None,
              line_w=0.5, corner_adj=None):
        from pptx.util import Inches, Pt
        shp = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
        if fill_hex:
            shp.fill.solid()
            shp.fill.fore_color.rgb = self._rgb(fill_hex)
        else:
            shp.fill.background()
        if line_hex:
            shp.line.color.rgb = self._rgb(line_hex)
            shp.line.width = Pt(line_w)
        else:
            shp.line.fill.background()
        if corner_adj is not None:
            self._round_corners(shp, corner_adj)
        return shp

    # ── Chrome ───────────────────────────────────────────────────────────

    def _bg_picture(self, slide, filename: str):
        from pptx.util import Inches
        path = os.path.join(_ASSETS_DIR, filename)
        slide.shapes.add_picture(path, Inches(0), Inches(0), Inches(SLIDE_W), Inches(SLIDE_H))

    def _chrome(self, slide, section_num: int, section_label: str, page_num: int | None):
        """Section eyebrow pill (top-left) + page number (bottom-right). Both
        omitted on cover/closing/static-intro slides, matching the sample."""
        from pptx.util import Inches
        label = f"SECTION {section_num} · {section_label}"
        pill_w = self._est_w(label, 11)
        self._rect(slide, PAD_X, 0.36, pill_w, 0.30, fill_hex=_C["blue"], corner_adj=500000)
        self._text(slide, PAD_X, 0.38, pill_w, 0.27, label, 11, _C["white"],
                   bold=True, align="CENTER", spc=80)
        if page_num is not None:
            self._text(slide, SLIDE_W - PAD_X - 0.60, 6.70, 0.60, 0.27,
                       str(page_num).zfill(2), 11, _C["gray"], align="RIGHT")

    def _title(self, slide, text: str):
        self._text(slide, PAD_X, TITLE_Y, CONTENT_W, TITLE_H, text, 30, _C["ink"],
                   bold=True, align="CENTER")

    def _card(self, slide, left, top, width, height, accent_hex, accent_h=0.07):
        """White card, subtle rounding, thin border, colored top accent stripe."""
        self._rect(slide, left, top, width, height, fill_hex=_C["white"],
                   line_hex=_C["border"], line_w=1.0, corner_adj=1744)
        self._rect(slide, left, top, width, accent_h, fill_hex=accent_hex)

    def _table_style(self, table, headers: list[str], rows: list[list[str]],
                      col_widths: list[float]):
        from pptx.util import Inches, Pt
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
        for c, h in enumerate(headers):
            cell = table.cell(0, c)
            cell.text = self._safe_text(h)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self._rgb(_C["navy"])
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.name = FONT
            run.font.size = Pt(11)
            run.font.bold = True
            run.font.color.rgb = self._rgb(_C["white"])
        for r, row in enumerate(rows, start=1):
            for c, val in enumerate(row):
                cell = table.cell(r, c)
                cell.text = self._safe_text(str(val))
                cell.fill.solid()
                cell.fill.fore_color.rgb = self._rgb(_C["white"])
                run = cell.text_frame.paragraphs[0].runs[0]
                run.font.name = FONT
                run.font.size = Pt(11)
                run.font.color.rgb = self._rgb(_C["ink"])

    # ── Cover / Agenda / Closing ─────────────────────────────────────────

    def _render_cover(self, slide, sd: dict):
        self._bg_picture(slide, "zbh-cover-bg.jpg")
        brand = sd.get("brand", "")
        industry = sd.get("industry", "")
        campaign_line = sd.get("campaign_line", "")
        date = sd.get("date", "")
        track = sd.get("track", "")
        prepared_by = sd.get("prepared_by", "AdtimaBox Sales Team")

        self._text(slide, PAD_X, 1.70, CONTENT_W, 0.30, brand.upper(), 12, _C["blue"],
                   bold=True, spc=160)
        self._text(slide, PAD_X, 2.15, CONTENT_W, 0.78, "ADTIMABOX PROPOSAL", 40,
                   _C["white"], bold=True)
        if campaign_line:
            self._text(slide, PAD_X, 3.05, CONTENT_W * 0.85, 0.35, campaign_line, 16,
                       _C["subtitle"])
        self._rect(slide, PAD_X, 3.75, 3.40, 0.05, fill_hex=_C["blue"])

        meta = [("DATE", date), ("INDUSTRY", industry), ("TRACK", track), ("PREPARED BY", prepared_by)]
        for i, (label, val) in enumerate(meta):
            x = PAD_X + i * 2.90
            self._text(slide, x, 6.05, 2.75, 0.27, label, 11, _C["blue"], bold=True, spc=100)
            self._text(slide, x, 6.33, 2.75, 0.30, val, 12, _C["white"])

        self._text(slide, PAD_X, 6.72, CONTENT_W, 0.27,
                   "Confidential. Prices exclude 8% VAT. Valid for 30 days from the issue date.",
                   11, _C["footnote"])

    def _render_agenda(self, slide, content_entries: list[tuple], n_static: int):
        self._bg_picture(slide, "zbh-frame.jpg")
        self._chrome_agenda(slide)
        self._title(slide, "What this proposal covers")

        # Compute slide-number ranges from the actual final layout: cover=1,
        # agenda=2, static intro=3..(2+n_static), then content slides sequential.
        static_start = 3
        static_end = static_start + n_static - 1
        next_page = static_end + 1

        ranges_by_section: dict[int, list[int]] = {}
        for ctype, section_num, section_label, _fallback in content_entries:
            ranges_by_section.setdefault(section_num, []).append(next_page)
            next_page += 1

        def _range_text(pages: list[int] | None) -> str:
            if not pages:
                return ""
            if len(pages) == 1:
                return f"Slide {pages[0]:02d}"
            return f"Slides {pages[0]:02d} to {pages[-1]:02d}"

        # Row height/fonts scale to however many sections exist (7 today) rather
        # than a fixed row count, so the agenda never overflows the slide if a
        # section is added later — sized to fill BODY_TOP..6.55" evenly, well
        # clear of the page-number chrome drawn at a fixed y=6.70 in
        # _chrome_agenda (a 7-row agenda measured at "end at 6.85" ran the last
        # row's bottom border straight through that page number).
        n_rows = max(len(_AGENDA_ITEMS), 1)
        row_gap = 0.08
        available_h = 6.55 - BODY_TOP
        row_h = max((available_h - row_gap * (n_rows - 1)) / n_rows, 0.45)
        badge_size = min(0.36, row_h - 0.20)
        title_size = 12 if row_h >= 0.65 else 10.5
        desc_size = 10 if row_h >= 0.65 else 9

        y = BODY_TOP
        for num, title, desc in _AGENDA_ITEMS:
            pages = ranges_by_section.get(num) if num != 1 else list(range(static_start, static_end + 1))
            self._rect(slide, PAD_X, y, CONTENT_W, row_h, fill_hex=_C["white"], corner_adj=6000)
            self._rect(slide, PAD_X, y, 0.05, row_h, fill_hex=_C["blue"])
            badge_y = y + (row_h - badge_size) / 2
            self._rect(slide, PAD_X + 0.26, badge_y, badge_size, badge_size, fill_hex=_C["blue"], corner_adj=50000)
            self._text(slide, PAD_X + 0.26, badge_y, badge_size, badge_size, f"{num:02d}", 9.5,
                       _C["white"], bold=True, align="CENTER")
            self._text(slide, PAD_X + 0.80, y + (row_h - 0.28) / 2, 2.90, 0.28, title,
                       title_size, _C["ink"], bold=True)
            self._text(slide, PAD_X + 3.80, y + (row_h - 0.42) / 2, 5.35, 0.42, desc,
                       desc_size, _C["gray"])
            self._text(slide, SLIDE_W - PAD_X - 1.85, y + (row_h - 0.27) / 2, 1.85, 0.27,
                       _range_text(pages), desc_size, _C["gray"], align="RIGHT")
            y += row_h + row_gap

    def _chrome_agenda(self, slide):
        """Agenda uses the same pill chrome as content slides but a fixed
        'AGENDA' label instead of a numbered section — matches the sample."""
        pill_w = self._est_w("AGENDA", 11)
        self._rect(slide, PAD_X, 0.36, pill_w, 0.30, fill_hex=_C["blue"], corner_adj=500000)
        self._text(slide, PAD_X, 0.38, pill_w, 0.27, "AGENDA", 11, _C["white"],
                   bold=True, align="CENTER", spc=80)
        self._text(slide, SLIDE_W - PAD_X - 0.60, 6.70, 0.60, 0.27, "02", 11,
                   _C["gray"], align="RIGHT")

    def _render_closing(self, slide, sd: dict):
        self._bg_picture(slide, "zbh-cover-bg.jpg")
        brand = sd.get("brand", "")
        date = sd.get("date", "")
        prepared_by = sd.get("prepared_by", "AdtimaBox")
        notes = sd.get("notes") or [
            "This proposal is confidential.",
            "Prices are valid for 30 days from the proposal date.",
        ]

        self._text(slide, PAD_X, 2.60, CONTENT_W, 0.30, "ADTIMABOX BY ADTIMA", 12,
                   _C["blue"], bold=True, spc=100)
        subtitle = f"Prepared by {prepared_by} | {date}" if date else f"Prepared by {prepared_by}"
        self._text(slide, PAD_X, 3.05, CONTENT_W, 0.66, subtitle, 24, _C["white"], bold=True)
        self._rect(slide, PAD_X, 3.95, 3.40, 0.05, fill_hex=_C["blue"])

        y = 4.25
        for note in notes[:2]:
            self._text(slide, PAD_X, y, CONTENT_W * 0.85, 0.33, note, 14, _C["subtitle"])
            y += 0.35

        self._text(slide, PAD_X, 6.72, CONTENT_W, 0.27,
                   "Confidential, for discussion between the two parties only.", 11,
                   _C["footnote"])

    # ── Executive summary (Section 1 of the source proposal) ────────────

    def _render_executive_summary(self, slide, sd: dict):
        headline = sd.get("headline", "")
        summary = sd.get("summary", "")
        metrics = (sd.get("metrics") or [])[:4]

        y = BODY_TOP
        if headline:
            h = self._est_h(headline, CONTENT_W, 16, bold=True) + 0.10
            self._text(slide, PAD_X, y, CONTENT_W, h, headline, 16, _C["blue"], bold=True)
            y += h + 0.16
        if summary:
            h = self._est_h(summary, CONTENT_W, 12) + 0.10
            self._text(slide, PAD_X, y, CONTENT_W, h, summary, 12, _C["ink"])
            y += h + 0.34

        if metrics:
            n = len(metrics)
            gap = 0.24
            card_w = (CONTENT_W - gap * (n - 1)) / n
            card_h = min(SLIDE_H - y - 0.85, 1.70)
            accents = [_C["blue"], _C["cyan"], _C["red"], _C["blue"]]
            value_h, label_gap, label_h = 0.55, 0.58, 0.45
            top_pad = max((card_h - (label_gap + label_h)) / 2, 0.10)
            for i, m in enumerate(metrics):
                x = PAD_X + i * (card_w + gap)
                self._card(slide, x, y, card_w, card_h, accents[i % len(accents)])
                self._text(slide, x + 0.14, y + top_pad, card_w - 0.28, value_h,
                           m.get("value", ""), 26, _C["ink"], bold=True, align="CENTER",
                           auto_fit=True)
                self._text(slide, x + 0.14, y + top_pad + label_gap, card_w - 0.28, label_h,
                           m.get("label", ""), 10, _C["gray"], align="CENTER")

    # ── Next steps (Section 7 of the source proposal — timeline + decisions) ──

    def _render_next_steps(self, slide, sd: dict):
        weeks = (sd.get("weeks") or [])[:4]
        decisions = (sd.get("decisions") or [])[:4]
        tech_items = (sd.get("tech_items") or [])[:4]

        y = BODY_TOP
        top_h = 2.20
        if weeks:
            self._card(slide, PAD_X, y, CONTENT_W, top_h, _C["blue"])
            self._text(slide, PAD_X + 0.26, y + 0.20, CONTENT_W - 0.52, 0.27,
                       "IMPLEMENTATION TIMELINE", 11, _C["gray"], bold=True)
            avail = top_h - 0.55
            row_h = avail / len(weeks)
            wy = y + 0.52
            for w in weeks:
                items_text = " · ".join(w.get("items") or [])
                label_line = f"{w.get('label', '')} — {items_text}" if items_text else w.get("label", "")
                self._text(slide, PAD_X + 0.26, wy, 1.35, row_h - 0.06, w.get("week", ""),
                           10, _C["blue"], bold=True)
                self._text(slide, PAD_X + 1.70, wy, CONTENT_W - 1.96, row_h - 0.06,
                           label_line, 10, _C["ink"])
                wy += row_h
            y += top_h + 0.18
        else:
            y += 0.0

        col_h = SLIDE_H - y - 0.85
        col_w = (CONTENT_W - 0.28) / 2

        def _col(x, label, items, accent):
            self._card(slide, x, y, col_w, col_h, accent)
            self._text(slide, x + 0.24, y + 0.20, col_w - 0.48, 0.27, label, 11,
                       _C["gray"], bold=True)
            avail = col_h - 0.55
            kept, heights, more = self._fit_list_height(items, col_w - 0.44, 10, avail, min_item_h=0.32)
            cy = y + 0.52
            for it, h in zip(kept, heights):
                text = it.get("text", "") if isinstance(it, dict) else it
                self._rect(slide, x + 0.26, cy + 0.06, 0.075, 0.075, fill_hex=accent, corner_adj=50000)
                self._text(slide, x + 0.46, cy, col_w - 0.72, h, text, 10, _C["ink"])
                cy += h
            if more:
                self._text(slide, x + 0.26, cy, col_w - 0.5, 0.20, f"+{more} more...", 9,
                           _C["gray"], italic=True)

        if decisions:
            _col(PAD_X, "KEY DECISIONS REQUIRED", decisions, _C["blue"])
        if tech_items:
            _col(PAD_X + col_w + 0.28, "TECH CONFIRMATION REQUIRED", tech_items, _C["cyan"])

    # ── Client requirements (2x2 grid) ───────────────────────────────────

    def _render_client_requirements(self, slide, sd: dict):
        quads = [
            ("CURRENT STATE", sd.get("current_state", ""), _C["red"]),
            ("CORE PAIN", sd.get("core_pain", ""), _C["red"]),
            ("DESIRED OUTCOME", sd.get("desired_outcome", ""), _C["blue"]),
            ("GAP", sd.get("gap", ""), _C["cyan"]),
        ]
        card_w = (CONTENT_W - 0.30) / 2
        card_h = 1.72
        row_gap = 0.18
        positions = [
            (PAD_X, BODY_TOP), (PAD_X + card_w + 0.30, BODY_TOP),
            (PAD_X, BODY_TOP + card_h + row_gap), (PAD_X + card_w + 0.30, BODY_TOP + card_h + row_gap),
        ]
        for (label, text, accent), (x, y) in zip(quads, positions):
            self._card(slide, x, y, card_w, card_h, accent)
            self._text(slide, x + 0.26, y + 0.26, card_w - 0.52, 0.27, label, 11,
                       accent, bold=True, spc=100)
            self._text(slide, x + 0.26, y + 0.60, card_w - 0.52, card_h - 0.80, text, 12, _C["ink"])

    # ── Solution: package ────────────────────────────────────────────────

    def _render_solution_package(self, slide, sd: dict):
        addons = (sd.get("addons") or [])[:8]
        package = sd.get("package") or {}
        tech_items = (sd.get("tech_confirm_items") or [])[:5]

        left_w = 7.40
        left_h = 3.90
        self._card(slide, PAD_X, BODY_TOP, left_w, left_h, _C["blue"])
        self._text(slide, PAD_X + 0.26, BODY_TOP + 0.24, left_w - 0.52, 0.27,
                   "CAMPAIGN INSTANT ADD-ONS", 11, _C["gray"], bold=True, spc=60)
        avail = left_h - 0.60
        kept, heights, more = self._fit_list_height(addons, left_w - 0.90, 10.5, avail, min_item_h=0.34)
        y = BODY_TOP + 0.56
        for item, h in zip(kept, heights):
            self._rect(slide, PAD_X + 0.28, y + 0.06, 0.075, 0.075, fill_hex=_C["blue"], corner_adj=50000)
            self._text(slide, PAD_X + 0.48, y, left_w - 0.75, h, item, 10.5, _C["ink"])
            y += h
        if more:
            self._text(slide, PAD_X + 0.28, y, left_w - 0.55, 0.22, f"+{more} more...", 9,
                       _C["gray"], italic=True)

        rx = PAD_X + left_w + 0.28
        right_w = CONTENT_W - left_w - 0.28

        pkg_h = 1.50
        self._card(slide, rx, BODY_TOP, right_w, pkg_h, _C["blue"])
        self._text(slide, rx + 0.26, BODY_TOP + 0.24, right_w - 0.52, 0.27, "PACKAGE", 11,
                   _C["gray"], bold=True, spc=60)
        self._text(slide, rx + 0.26, BODY_TOP + 0.55, right_w - 0.52, 0.38,
                   package.get("name", ""), 16, _C["ink"], bold=True)
        self._text(slide, rx + 0.26, BODY_TOP + 0.95, right_w - 0.52, 0.55,
                   package.get("tier_note", ""), 10, _C["gray"])

        tech_y = BODY_TOP + pkg_h + 0.18
        tech_h = left_h - pkg_h - 0.18
        self._card(slide, rx, tech_y, right_w, tech_h, _C["cyan"])
        label_h = self._est_h("CUSTOM ITEMS REQUIRING TECH CONFIRMATION", right_w - 0.52, 11, bold=True) + 0.06
        self._text(slide, rx + 0.26, tech_y + 0.24, right_w - 0.52, label_h,
                   "CUSTOM ITEMS REQUIRING TECH CONFIRMATION", 11, _C["gray"], bold=True, spc=60)
        list_start = tech_y + 0.30 + label_h
        t_avail = tech_h - (list_start - tech_y) - 0.10
        t_kept, t_heights, t_more = self._fit_list_height(
            tech_items, right_w - 0.55, 10, t_avail, min_item_h=0.34)
        ty = list_start
        for item, h in zip(t_kept, t_heights):
            self._rect(slide, rx + 0.28, ty + 0.06, 0.075, 0.075, fill_hex=_C["cyan"], corner_adj=50000)
            self._text(slide, rx + 0.48, ty, right_w - 0.75, h, item, 10, _C["ink"])
            ty += h
        if t_more:
            self._text(slide, rx + 0.28, ty, right_w - 0.55, 0.22, f"+{t_more} more...", 9,
                       _C["gray"], italic=True)

    # ── Solution: user journey (icon-flow) ───────────────────────────────

    def _render_user_journey(self, slide, sd: dict):
        steps = (sd.get("steps") or [])[:5]
        footer = sd.get("footer", "")
        n = max(len(steps), 1)
        gap = 0.27
        card_w = (CONTENT_W - gap * (n - 1)) / n
        card_h = 2.40
        y = 2.10
        role_colors = {"consumer": _C["blue"], "system": _C["cyan"], "admin": _C["red"], "staff": _C["gray"]}

        for i, st in enumerate(steps):
            x = PAD_X + i * (card_w + gap)
            role = st.get("role", "consumer")
            rc = role_colors.get(role, _C["blue"])
            self._card(slide, x, y, card_w, card_h, rc, accent_h=0.06)
            self._text(slide, x + 0.18, y + 0.22, card_w - 0.36, 0.27, role.upper(), 9,
                       _C["gray"], bold=True)
            badge_y = y + 0.62
            self._rect(slide, x + 0.18, badge_y, 0.40, 0.40, fill_hex=rc, corner_adj=50000)
            self._text(slide, x + 0.18, badge_y, 0.40, 0.40, str(st.get("number", i + 1)),
                       11, _C["white"], bold=True, align="CENTER")
            self._text(slide, x + 0.18, y + 1.18, card_w - 0.36, 0.28, st.get("label", ""),
                       11, _C["ink"], bold=True)
            self._text(slide, x + 0.18, y + 1.52, card_w - 0.36, 0.63, st.get("desc", ""),
                       9.5, _C["gray"])

            if i < n - 1:
                arr_x = x + card_w + 0.03
                arr_w = max(gap - 0.06, 0.05)
                self._rect(slide, arr_x, y + card_h * 0.32, arr_w, 0.015, fill_hex=_C["gray"])

        if footer:
            fy = y + card_h + 0.22
            fh = 0.90
            self._card(slide, PAD_X, fy, CONTENT_W, fh, _C["blue"])
            self._text(slide, PAD_X + 0.30, fy + 0.18, CONTENT_W - 0.60, 0.27, "JOURNEY", 11,
                       _C["gray"], bold=True)
            self._text(slide, PAD_X + 0.30, fy + 0.50, CONTENT_W - 0.60, fh - 0.55, footer, 10.5,
                       _C["ink"])

    # ── Solution: flowchart (new layout, no precedent) ───────────────────

    def _render_solution_flowchart(self, slide, sd: dict):
        nodes = (sd.get("nodes") or [])[:8]
        side_note = sd.get("side_note", "")
        if not nodes:
            return

        rows: list[list[dict]] = []
        for i in range(0, len(nodes), 4):
            rows.append(nodes[i:i + 4])

        row_h = 0.85
        row_gap = 0.45
        y = 2.15
        for row in rows:
            n = len(row)
            gap = 0.30
            box_w = (CONTENT_W - gap * (n - 1)) / n
            x = PAD_X
            for node in row:
                is_decision = bool(node.get("decision"))
                text = node.get("text", "")
                if is_decision:
                    self._rect(slide, x, y, box_w, row_h, fill_hex=_C["white"],
                               line_hex=_C["blue"], line_w=1.25, corner_adj=50000)
                    self._text(slide, x + 0.12, y, box_w - 0.24, row_h, text, 11,
                               _C["ink"], bold=True, align="CENTER")
                else:
                    self._card(slide, x, y, box_w, row_h, _C["blue"], accent_h=0.05)
                    self._text(slide, x + 0.18, y + 0.22, box_w - 0.36, row_h - 0.30,
                               text, 10.5, _C["ink"])
                if x + box_w < PAD_X + CONTENT_W - 0.05:
                    arr_x = x + box_w + 0.03
                    self._rect(slide, arr_x, y + row_h / 2 - 0.01, max(gap - 0.06, 0.05), 0.015,
                               fill_hex=_C["gray"])
                x += box_w + gap
            y += row_h + row_gap

        if side_note:
            self._text(slide, PAD_X, y, CONTENT_W, 0.27, side_note, 10, _C["gray"], italic=True)

    # ── Solution: touchpoints table ──────────────────────────────────────

    def _render_touchpoints_table(self, slide, sd: dict):
        rows = (sd.get("rows") or [])[:6]
        note = sd.get("note", "")
        if not rows:
            return
        headers = ["Timing", "Message content", "Message type", "Channel"]
        col_w = [2.40, 4.60, 2.40, 2.53]
        n_rows = len(rows) + 1
        row_h = 0.58
        from pptx.util import Inches
        gfx = slide.shapes.add_table(n_rows, len(headers), Inches(PAD_X), Inches(BODY_TOP),
                                     Inches(CONTENT_W), Inches(row_h * n_rows))
        table = gfx.table
        data_rows = [[r.get("timing", ""), r.get("message", ""), r.get("type", ""),
                      r.get("channel", "")] for r in rows]
        self._table_style(table, headers, data_rows, col_w)

        if note:
            ny = BODY_TOP + row_h * n_rows + 0.20
            self._card(slide, PAD_X, ny, CONTENT_W, 0.92, _C["blue"])
            self._text(slide, PAD_X + 0.26, ny + 0.20, CONTENT_W - 0.52, 0.60, note, 11, _C["ink"])

    # ── Solution: compliance ──────────────────────────────────────────────

    def _render_compliance(self, slide, sd: dict):
        verdict = (sd.get("verdict") or "CONDITIONS").upper()
        verdict_label = {
            "CLEAR": "CLEAR TO PROCEED",
            "CONDITIONS": "PROCEED WITH CONDITIONS",
            "BLOCKED": "BLOCKED",
        }.get(verdict, verdict)
        verdict_colors = {"CLEAR": _C["blue"], "CONDITIONS": "C8932B", "BLOCKED": _C["red"]}
        vc = verdict_colors.get(verdict, _C["blue"])
        conditions = (sd.get("conditions") or [])[:4]
        docs = (sd.get("docs_required") or [])[:4]
        consent_text = sd.get("consent_text", "")

        pill_w = self._est_w(verdict_label, 12, pad_in=0.30)
        self._rect(slide, PAD_X, BODY_TOP, pill_w, 0.27, fill_hex=vc, corner_adj=200000)
        self._text(slide, PAD_X, BODY_TOP, pill_w, 0.27, verdict_label, 11, _C["white"],
                   bold=True, align="CENTER")

        col_y = BODY_TOP + 0.46
        col_h = 2.30
        col_w = (CONTENT_W - 0.28) / 2

        def _col(x, label, items):
            self._card(slide, x, col_y, col_w, col_h, _C["blue"])
            self._text(slide, x + 0.24, col_y + 0.22, col_w - 0.48, 0.27, label, 11,
                       _C["gray"], bold=True)
            avail = col_h - 0.55
            kept, heights, more = self._fit_list_height(items, col_w - 0.44, 10, avail, min_item_h=0.34)
            cy = col_y + 0.58
            for it, h in zip(kept, heights):
                self._rect(slide, x + 0.26, cy + 0.06, 0.075, 0.075, fill_hex=_C["blue"], corner_adj=50000)
                self._text(slide, x + 0.46, cy, col_w - 0.72, h, it, 10, _C["ink"])
                cy += h
            if more:
                self._text(slide, x + 0.26, cy, col_w - 0.5, 0.20, f"+{more} more...", 9,
                           _C["gray"], italic=True)

        if conditions:
            _col(PAD_X, "CONDITIONS BEFORE LAUNCH", conditions)
        if docs:
            _col(PAD_X + col_w + 0.28, "DOCUMENTS THE CLIENT PROVIDES", docs)

        if consent_text:
            cy = col_y + col_h + 0.18
            self._card(slide, PAD_X, cy, CONTENT_W, 1.00, _C["blue"])
            self._text(slide, PAD_X + 0.26, cy + 0.18, CONTENT_W - 0.52, 0.27,
                       "MANDATORY CONSENT TEXT", 11, _C["gray"], bold=True)
            self._text(slide, PAD_X + 0.26, cy + 0.50, CONTENT_W - 0.52, 0.45, consent_text, 10.5,
                       _C["ink"])

    # ── Quotation ─────────────────────────────────────────────────────────

    def _quotation_table_and_card(self, slide, sd: dict, when_to_choose: list | None = None):
        line_items = (sd.get("line_items") or [])[:6]
        subtotal = sd.get("subtotal", "")
        vat = sd.get("vat", "")
        total = sd.get("total", "")
        notes = (sd.get("notes") or [])[:4]

        if line_items:
            n_rows = len(line_items) + 1
            row_h = 0.58
            from pptx.util import Inches
            gfx = slide.shapes.add_table(n_rows, 2, Inches(PAD_X), Inches(BODY_TOP),
                                         Inches(7.00), Inches(row_h * n_rows))
            table = gfx.table
            data_rows = [[li.get("label", ""), li.get("amount", "")] for li in line_items]
            self._table_style(table, ["Line item", "Amount (VND)"], data_rows, [4.60, 2.40])

        rx = PAD_X + 7.00 + 0.28
        right_w = CONTENT_W - 7.00 - 0.28

        breakdown = [("Subtotal excl. VAT", subtotal), ("VAT 8%", vat)]
        card_h = 0.33 * len(breakdown) + 0.60
        self._card(slide, rx, BODY_TOP, right_w, card_h, _C["blue"])
        y = BODY_TOP + 0.24
        for label, val in breakdown:
            self._text(slide, rx + 0.24, y, right_w * 0.55, 0.27, label, 10.5, _C["ink"])
            self._text(slide, rx + right_w * 0.55, y, right_w * 0.42, 0.27, val, 10.5,
                       _C["ink"], align="RIGHT")
            y += 0.33
        self._rect(slide, rx + 0.24, y + 0.04, right_w - 0.48, 0.32, fill_hex=_C["navy"], corner_adj=8000)
        self._text(slide, rx + 0.30, y + 0.08, right_w * 0.5, 0.24, "TOTAL", 10.5, _C["white"], bold=True)
        self._text(slide, rx + right_w * 0.50, y + 0.08, right_w * 0.44, 0.24, total, 10.5,
                   _C["white"], bold=True, align="RIGHT")

        notes_y = BODY_TOP + card_h + 0.18
        label = "WHEN TO CHOOSE THIS OPTION" if when_to_choose else "BUDGET NOTES"
        items = when_to_choose if when_to_choose else notes
        if items:
            notes_h = SLIDE_H - notes_y - 0.85
            self._card(slide, rx, notes_y, right_w, notes_h, _C["cyan"])
            self._text(slide, rx + 0.24, notes_y + 0.20, right_w - 0.48, 0.27, label, 10.5,
                       _C["gray"], bold=True)
            kept, heights, more = self._fit_list_height(
                items, right_w - 0.55, 9.5, notes_h - 0.55, min_item_h=0.30)
            cy = notes_y + 0.52
            for it, h in zip(kept, heights):
                self._rect(slide, rx + 0.26, cy + 0.06, 0.07, 0.07, fill_hex=_C["cyan"], corner_adj=50000)
                self._text(slide, rx + 0.44, cy, right_w - 0.70, h, it, 9.5, _C["ink"])
                cy += h

    def _render_quotation(self, slide, sd: dict):
        self._quotation_table_and_card(slide, sd)
        recon = sd.get("reconciliation_note", "")
        if recon:
            self._text(slide, PAD_X, 5.36, 7.00, 0.60, recon, 10, _C["gray"])

    def _render_quotation_alternative(self, slide, sd: dict):
        self._quotation_table_and_card(slide, sd, when_to_choose=sd.get("when_to_choose") or [])
        desc = sd.get("description", "")
        if desc:
            self._text(slide, PAD_X, 3.75, 7.00, 0.60, desc, 10.5, _C["gray"])

    # ── Case study ────────────────────────────────────────────────────────

    def _render_case_study(self, slide, sd: dict):
        cases = (sd.get("cases") or [])[:2]
        disclaimer = sd.get("disclaimer", "")
        if not cases:
            return
        gap = 0.30
        card_w = (CONTENT_W - gap) / 2
        card_h = 3.55

        rows = [
            ("WHY IT IS RELEVANT", "why_relevant"),
            ("WHAT WAS DONE", "what_was_done"),
            ("RESULT", "result"),
            ("APPLICABLE TO THE CLIENT BECAUSE", "applies_because"),
        ]

        for i, case in enumerate(cases):
            x = PAD_X + i * (card_w + gap)
            self._card(slide, x, BODY_TOP, card_w, card_h, _C["blue"])
            badge_w = 0.90
            self._rect(slide, x + 0.26, BODY_TOP + 0.26, badge_w, 0.27, fill_hex=_C["blue"], corner_adj=100000)
            self._text(slide, x + 0.26, BODY_TOP + 0.26, badge_w, 0.27, case.get("alias", ""),
                       10, _C["white"], bold=True, align="CENTER")
            self._text(slide, x + 0.26 + badge_w + 0.16, BODY_TOP + 0.22, card_w - badge_w - 0.60,
                       0.30, case.get("name", ""), 12, _C["ink"], bold=True)

            ry = BODY_TOP + 0.70
            for label, key in rows:
                self._text(slide, x + 0.26, ry, card_w - 0.52, 0.27, label, 11, _C["blue"],
                           bold=True)
                self._text(slide, x + 0.26, ry + 0.24, card_w - 0.52, 0.36, case.get(key, ""),
                           10.5, _C["ink"])
                ry += 0.69

        if disclaimer:
            self._text(slide, PAD_X, BODY_TOP + card_h + 0.18, CONTENT_W, 0.27, disclaimer, 10,
                       _C["gray"], italic=True)


def create_corporate_pptx_generator() -> CorporatePPTXGenerator:
    return CorporatePPTXGenerator()
