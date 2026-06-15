# PPTX Generation Module
# ======================
# Generates PowerPoint presentations from approved plans/proposals.
# Uses python-pptx to open a branded master template and fill placeholders.

import io
import os
from pathlib import Path
from typing import Optional
from datetime import datetime


class PPTXGenerator:
    """Generates PowerPoint presentations from approved plans/proposals."""

    def __init__(self, template_path: Optional[str] = None):
        """
        Initialize the PPTX generator.

        Args:
            template_path: Path to a branded master .pptx template.
                           If None, creates a basic presentation.
        """
        self.template_path = template_path

    async def generate(
        self,
        plan_data: dict,
        client_name: str = "Client",
        output_path: Optional[str] = None
    ) -> dict:
        """
        Generate a PPTX deck from plan data.

        Args:
            plan_data: The approved plan/proposal data containing:
                - title: Project title
                - summary: Executive summary
                - items: Line items/services
                - total: Total amount
                - timeline: Project timeline
                - solutions: Proposed solutions
            client_name: Name of the client
            output_path: Optional path to save the file

        Returns:
            dict with status, file_path, and preview info
        """
        try:
            # Try to use python-pptx if available
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor

            prs = Presentation()

            # Title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]

            title.text = plan_data.get("title", "Sales Proposal")
            subtitle.text = f"{client_name}\n{datetime.now().strftime('%B %Y')}"

            # Agenda/Overview slide
            agenda_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(agenda_layout)
            title = slide.shapes.title
            title.text = "Proposal Overview"

            # Add content
            content = slide.placeholders[1]
            tf = content.text_frame
            tf.text = plan_data.get("summary", "Summary of proposed solution")

            # Solutions slide
            if "solutions" in plan_data:
                bullet_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(bullet_layout)
                title = slide.shapes.title
                title.text = "Proposed Solutions"

                content = slide.placeholders[1]
                tf = content.text_frame
                tf.text = "Key Solutions:"

                for solution in plan_data.get("solutions", []):
                    p = tf.add_paragraph()
                    p.text = f"• {solution}"
                    p.level = 1

            # Timeline slide
            if "timeline" in plan_data:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                title = slide.shapes.title
                title.text = "Project Timeline"

                content = slide.placeholders[1]
                tf = content.text_frame
                tf.text = plan_data.get("timeline", "Timeline overview")

            # Pricing slide
            if "items" in plan_data or "total" in plan_data:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                title = slide.shapes.title
                title.text = "Investment"

                content = slide.placeholders[1]
                tf = content.text_frame

                if "items" in plan_data:
                    for item in plan_data.get("items", []):
                        p = tf.add_paragraph()
                        item_text = item.get("name", "Item")
                        if "price" in item:
                            item_text += f": {item['price']}"
                        p.text = item_text
                        p.level = 1

                if "total" in plan_data:
                    p = tf.add_paragraph()
                    p.text = f"Total: {plan_data['total']}"
                    p.level = 0
                    p.font.bold = True

            # Next Steps slide
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            title.text = "Next Steps"

            content = slide.placeholders[1]
            tf = content.text_frame
            tf.text = "1. Review and approval\n2. Contract signing\n3. Project kickoff\n4. Implementation"

            # Save to buffer
            buffer = io.BytesIO()
            prs.save(buffer)
            buffer.seek(0)

            # Save to file if path provided
            file_path = None
            if output_path:
                os.makedirs(os.path.dirname(output_path), exist_ok=True)
                with open(output_path, 'wb') as f:
                    f.write(buffer.getvalue())
                file_path = output_path

            return {
                "status": "success",
                "file_path": file_path,
                "preview": {
                    "slides": len(prs.slides),
                    "title": plan_data.get("title", "Sales Proposal"),
                    "has_timeline": "timeline" in plan_data,
                    "has_pricing": "total" in plan_data
                },
                "buffer": buffer.getvalue() if not output_path else None
            }

        except ImportError:
            # Fallback: create a simple text-based preview if python-pptx not available
            return await self._generate_fallback(plan_data, client_name)

    async def _generate_fallback(
        self,
        plan_data: dict,
        client_name: str
    ) -> dict:
        """Generate a simple text preview when python-pptx is not available."""

        content = f"""
# {plan_data.get('title', 'Sales Proposal')}

## Client: {client_name}

## Overview
{plan_data.get('summary', 'No summary available')}

## Solutions
"""
        for solution in plan_data.get("solutions", []):
            content += f"- {solution}\n"

        if "timeline" in plan_data:
            content += f"\n## Timeline\n{plan_data.get('timeline')}\n"

        if "total" in plan_data:
            content += f"\n## Investment\nTotal: {plan_data.get('total')}\n"

        content += "\n## Next Steps\n1. Review and approval\n2. Contract signing\n3. Project kickoff\n4. Implementation"

        return {
            "status": "success",
            "preview": content,
            "slides": 5,
            "file_path": None,
            "fallback": True
        }

    def get_supported_features(self) -> list[str]:
        """Return list of supported PPTX features."""
        return [
            "title_slide",
            "agenda_slide",
            "solutions_slide",
            "timeline_slide",
            "pricing_slide",
            "next_steps_slide"
        ]


# Factory function for creating generator
def create_pptx_generator(template_path: Optional[str] = None) -> PPTXGenerator:
    """Create a PPTX generator instance."""
    return PPTXGenerator(template_path=template_path)