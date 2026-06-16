"""
Design Agent - Generates PPTX and Figma artifacts
"""

from typing import Optional
from pathlib import Path

from schemas.state import AgentOutput, SalesCaseState
from agents.base import BaseAgent

# Import design backend - use relative import
from pathlib import Path

# Add project root to path for design backend
import sys
import os
_agent_dir = Path(__file__).parent  # agents/design/
_backend_dir = _agent_dir.parent  # backend/
sys.path.insert(0, str(_backend_dir))

from design.backend import create_design_backend


def get_design_agent() -> BaseAgent:
    """Get the design agent instance."""
    return DesignAgent()


class DesignAgent(BaseAgent):
    """Agent that generates design artifacts: PPTX and Figma wireframes."""

    def __init__(self):
        super().__init__(
            name="design",
            model_key="MODEL_DESIGN",
            role_description="Design & Slide Creator - generates PPTX presentations and Figma wireframes",
            prompt_path="backend/agents/design/prompt.md",
            knowledge_dir="backend/agents/design/knowledge",
            skills_dir="backend/agents/design/skills",
        )
        # Initialize design backend
        self.design_backend = create_design_backend()

    async def run(self, state: SalesCaseState) -> AgentOutput:
        """Run the design agent to generate artifacts."""

        # Get context from other agents
        market_output = state.outputs.get("market_strategy")
        product_output = state.outputs.get("product_solution")

        # Build context for artifact generation
        context = self._build_context(state, market_output, product_output)

        # Generate PPTX
        pptx_path = await self._generate_pptx(context, state.session_id)

        # Generate Figma/FigJam wireframe
        figma_result = await self._generate_figma_wireframe(context, state.session_id)

        # Build response
        artifacts = []
        if pptx_path:
            artifacts.append({
                "type": "pptx",
                "path": pptx_path,
                "description": "Presentation deck"
            })

        if figma_result and figma_result.get("status") == "success":
            artifacts.append({
                "type": "figma",
                "data": figma_result,
                "description": "Wireframe design"
            })

        summary = f"Generated design artifacts: PPTX presentation"
        if figma_result and figma_result.get("status") == "success":
            summary += ", Figma/FigJam wireframe"

        return AgentOutput(
            agent="design",
            status="COMPLETE",
            payload={
                "artifacts": artifacts,
                "pptx_path": pptx_path,
                "figma_result": figma_result,
            },
            summary=summary,
            confidence=0.9,
            needs=None,
            questions=[],
        )

    def _build_context(self, state, market_output, product_output) -> dict:
        """Build context from other agent outputs."""
        context = {
            "session_id": state.session_id,
            "brief": state.brief.model_dump() if state.brief else {},
        }

        if market_output and hasattr(market_output, "payload"):
            context["market_strategy"] = market_output.payload

        if product_output and hasattr(product_output, "payload"):
            context["product_solution"] = product_output.payload

        return context

    async def _generate_pptx(self, context: dict, session_id: str) -> Optional[str]:
        """Generate PPTX file."""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from datetime import datetime

            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            brief = context.get("brief", {})
            brand = brief.get("brand", "Brand")
            goal = brief.get("goal", "Campaign")

            # Slide 1: Title
            title_slide = prs.slides.add_slide(prs.slide_layouts[6])
            title = title_slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1.5))
            tf = title.text_frame
            p = tf.paragraphs[0]
            p.text = f"{brand} - Campaign Proposal"
            p.font.size = Pt(44)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

            # Slide 2: Overview
            overview_slide = prs.slides.add_slide(prs.slide_layouts[6])
            title_box = overview_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
            tf = title_box.text_frame
            tf.paragraphs[0].text = "Campaign Overview"
            tf.paragraphs[0].font.size = Pt(32)
            tf.paragraphs[0].font.bold = True

            content_box = overview_slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(5))
            tf = content_box.text_frame
            tf.word_wrap = True

            if brief:
                for key, value in brief.items():
                    if value:
                        p = tf.add_paragraph()
                        p.text = f"{key.replace('_', ' ').title()}: {value}"
                        p.font.size = Pt(18)
                        p.space_after = Pt(12)

            # Slide 3: Strategy (from market_strategy)
            market = context.get("market_strategy", {})
            if market:
                strategy_slide = prs.slides.add_slide(prs.slide_layouts[6])
                title_box = strategy_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
                tf = title_box.text_frame
                tf.paragraphs[0].text = "Market Strategy"
                tf.paragraphs[0].font.size = Pt(32)
                tf.paragraphs[0].font.bold = True

                content_box = strategy_slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(5))
                tf = content_box.text_frame
                tf.word_wrap = True

                # Extract summary from market strategy
                if isinstance(market, dict):
                    for key, value in list(market.items())[:5]:
                        p = tf.add_paragraph()
                        p.text = f"{key}: {str(value)[:200]}"
                        p.font.size = Pt(16)
                        p.space_after = Pt(8)

            # Slide 4: Product Solution
            product = context.get("product_solution", {})
            if product:
                product_slide = prs.slides.add_slide(prs.slide_layouts[6])
                title_box = product_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
                tf = title_box.text_frame
                tf.paragraphs[0].text = "Product Solution"
                tf.paragraphs[0].font.size = Pt(32)
                tf.paragraphs[0].font.bold = True

                content_box = product_slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(5))
                tf = content_box.text_frame
                tf.word_wrap = True
                for key, value in list(product.items())[:8]:
                    p = tf.add_paragraph()
                    p.text = f"{key}: {value}"
                    p.font.size = Pt(16)
                    p.space_after = Pt(8)

            # Slide 5: Pricing
            pricing = (context.get("product_solution", {}) or {}).get("pricing_breakdown", {})
            if pricing:
                pricing_slide = prs.slides.add_slide(prs.slide_layouts[6])
                title_box = pricing_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
                tf = title_box.text_frame
                tf.paragraphs[0].text = "Pricing & Timeline"
                tf.paragraphs[0].font.size = Pt(32)
                tf.paragraphs[0].font.bold = True

                content_box = pricing_slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(5))
                tf = content_box.text_frame
                tf.word_wrap = True

                for key, value in pricing.items():
                    p = tf.add_paragraph()
                    p.text = f"{key}: {value}"
                    p.font.size = Pt(18)
                    p.space_after = Pt(8)

            # Save
            output_dir = Path(__file__).parent.parent / "output"
            output_dir.mkdir(parents=True, exist_ok=True)
            output_dir.mkdir(exist_ok=True)
            output_path = output_dir / f"proposal_{session_id}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            prs.save(str(output_path))

            return str(output_path)

        except ImportError:
            # python-pptx not installed
            return None
        except Exception as e:
            print(f"Error generating PPTX: {e}")
            return None

    async def _generate_figma_description(self, context: dict) -> Optional[str]:
        """Generate Figma wireframe description."""
        # This returns a description that could be used to create a Figma wireframe
        brief = context.get("brief", {})

        description = f"""
## Figma Wireframe Concept

### Page 1: Landing / Home
- Hero section with brand logo and campaign tagline
- CTA button "Join Now" / "Quét QR"
- Brief instructions on how to participate

### Page 2: Game / Activity
- Interactive game area (spin wheel / scratch card / quiz)
- User progress indicator
- Points/voucher display

### Page 3: Data Entry Form
- Phone number input (required)
- Name input
- Optional: Email, birthday
- Consent checkbox for marketing

### Page 4: Voucher / Reward
- Voucher code display
- QR code for redemption
- Share buttons (Zalo, Facebook)
- "Redeem Now" CTA
"""
        return description
